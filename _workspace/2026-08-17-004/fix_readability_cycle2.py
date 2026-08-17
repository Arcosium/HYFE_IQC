from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt


PPTX = Path(
    "/home/arcosium/projects/GenomicWQB/docs/머신발표/GenomicWQB_머신발표.pptx"
)


def shape(slide, name: str):
    return next(item for item in slide.shapes if item.name == name)


def set_text(shape_, paragraphs: list[str], size: float) -> None:
    frame = shape_.text_frame
    current = list(frame.paragraphs)
    while len(current) < len(paragraphs):
        current.append(frame.add_paragraph())
    for paragraph, text in zip(current, paragraphs):
        if paragraph.runs:
            paragraph.runs[0].text = text
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = text
        for run in paragraph.runs:
            run.font.size = Pt(size)
    for paragraph in current[len(paragraphs) :]:
        frame._element.remove(paragraph._element)
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.word_wrap = True


def set_size(shape_, size: float) -> None:
    for paragraph in shape_.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)
    shape_.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    shape_.text_frame.word_wrap = True


def normalize_missing_sizes(presentation: Presentation) -> None:
    """LibreOffice must not fall back to a large theme size on later paragraphs."""
    for slide in presentation.slides:
        for item in slide.shapes:
            if not getattr(item, "has_text_frame", False):
                continue
            explicit = [
                run.font.size
                for paragraph in item.text_frame.paragraphs
                for run in paragraph.runs
                if run.text and run.font.size is not None
            ]
            if explicit:
                fallback = explicit[0]
                for paragraph in item.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text and run.font.size is None:
                            run.font.size = fallback
            item.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            item.text_frame.word_wrap = True


def main() -> None:
    presentation = Presentation(PPTX)
    normalize_missing_sizes(presentation)

    slide9 = presentation.slides[8]
    set_size(shape(slide9, "TextBox 22"), 9.2)
    set_size(shape(slide9, "TextBox 27"), 9.2)

    slide11 = presentation.slides[10]
    set_text(
        shape(slide11, "TextBox 15"),
        [
            "고회전 관문 미달은 FAIL이 아니라 WARNING이라 fail_items에 잡히지 않았습니다.",
            "부모의 실측 지표를 변이 지시에 더해 숨은 병목을 보완했습니다.",
        ],
        10.5,
    )
    set_text(
        shape(slide11, "TextBox 17"),
        [
            "누적 관측에서 boost 36.0%, smooth 0.5%로 약 72배 차이가 났습니다.",
            "통제 실험은 아니며 이 성공률이 다음 변이의 우선순위를 바꿉니다.",
        ],
        11.0,
    )

    slide12 = presentation.slides[11]
    set_text(
        shape(slide12, "TextBox 13"),
        [
            "· Sharpe(신호의 세기)",
            "· Fitness(회전율 대비 효율)",
            "· 제출 경로 점수(통과 근접도)",
            "· self-correlation(기존 제출작과의 거리)",
            "· 최근 2년 Sharpe(시간 안정성)",
        ],
        10.0,
    )

    slide13 = presentation.slides[12]
    set_text(
        shape(slide13, "TextBox 14"),
        ["완전 통과 보상은 희소하다 (35,701건 중 분석 조건상 442건, 1.2%)"],
        9.5,
    )
    for name in [
        "TextBox 13",
        "TextBox 15",
        "TextBox 18",
        "TextBox 19",
        "TextBox 20",
        "TextBox 22",
        "TextBox 23",
        "TextBox 24",
        "TextBox 27",
        "TextBox 28",
        "TextBox 29",
    ]:
        set_size(shape(slide13, name), 9.5)

    slide14 = presentation.slides[13]
    set_size(shape(slide14, "TextBox 13"), 10.5)
    set_size(shape(slide14, "TextBox 18"), 10.5)
    set_size(shape(slide14, "TextBox 23"), 10.5)
    set_size(shape(slide14, "TextBox 26"), 10.5)

    slide16 = presentation.slides[15]
    set_text(
        shape(slide16, "TextBox 15"),
        [
            "누적 빈도와 현재 병목은 다릅니다.",
            "8월 13~14일 태그 134건 중 53건의 prod 상관 중앙값은 0.9028이었고 0.7 미만은 없었습니다.",
            "약한 신호와 높은 유사도가 함께 만든 병목입니다.",
        ],
        10.5,
    )

    slide18 = presentation.slides[17]
    table_text = {
        "TextBox 24": "미달 피라미드 칸 우선",
        "TextBox 25": "3건이면 달성. 찬 칸은 상관만 높인다",
        "TextBox 31": "데이터셋이 섞이면 risk 칸으로 집계",
        "TextBox 36": "상한이 없으면 한 데이터셋으로 쏠린다",
        "TextBox 41": "prod 0.70~0.80을 중립화해 재시뮬",
        "TextBox 42": "0.8 이상은 회복 사례 없음",
    }
    for name, text in table_text.items():
        item = shape(slide18, name)
        set_text(item, [text], 10.0)
    for name in [
        "TextBox 22",
        "TextBox 23",
        "TextBox 28",
        "TextBox 29",
        "TextBox 30",
        "TextBox 33",
        "TextBox 34",
        "TextBox 35",
        "TextBox 39",
        "TextBox 40",
    ]:
        set_size(shape(slide18, name), 10.0)
    for name in [
        "TextBox 22",
        "TextBox 23",
        "TextBox 24",
        "TextBox 25",
        "TextBox 28",
        "TextBox 29",
        "TextBox 30",
        "TextBox 31",
        "TextBox 33",
        "TextBox 34",
        "TextBox 35",
        "TextBox 36",
        "TextBox 39",
        "TextBox 40",
        "TextBox 41",
        "TextBox 42",
    ]:
        shape(slide18, name).height = Inches(0.26)

    slide19 = presentation.slides[18]
    set_text(
        shape(slide19, "TextBox 13"),
        ["GLB/D1 Broker·Imbalance는 필드가 0이라 목표에서 제외했습니다."],
        10.5,
    )
    set_size(shape(slide19, "TextBox 14"), 11.0)

    slide20 = presentation.slides[19]
    set_size(shape(slide20, "TextBox 4"), 23.0)
    set_size(shape(slide20, "TextBox 6"), 15.0)

    presentation.save(PPTX)


if __name__ == "__main__":
    main()
