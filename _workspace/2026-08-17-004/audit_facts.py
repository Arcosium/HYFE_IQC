from __future__ import annotations

import datetime as dt
import json
import sqlite3
from pathlib import Path
from zoneinfo import ZoneInfo

from docx import Document
from pptx import Presentation


ROOT = Path("/home/arcosium/projects/GenomicWQB")
DB = ROOT / "data/hyfe_iqc.db"
CUTOFF = dt.datetime(2026, 8, 14, 16, 20, tzinfo=ZoneInfo("Asia/Seoul")).timestamp()


def scalar(conn: sqlite3.Connection, sql: str, params=()):
    return conn.execute(sql, params).fetchone()[0]


def ppt_text(path: Path) -> str:
    presentation = Presentation(path)
    return "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def docx_text(path: Path) -> str:
    document = Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


def main() -> None:
    conn = sqlite3.connect(DB)
    checks = {
        "alphas": (
            scalar(conn, "SELECT count(*) FROM alphas WHERE ts<=?", (CUTOFF,)),
            35_701,
        ),
        "rounds": (
            scalar(conn, "SELECT count(*) FROM rounds WHERE started_at<=?", (CUTOFF,)),
            923,
        ),
        "unique_code": (
            scalar(conn, "SELECT count(DISTINCT code) FROM alphas WHERE ts<=?", (CUTOFF,)),
            26_404,
        ),
        "unique_settings_fp": (
            scalar(
                conn,
                "SELECT count(DISTINCT settings_fp) FROM alphas "
                "WHERE ts<=? AND settings_fp IS NOT NULL",
                (CUTOFF,),
            ),
            2_227,
        ),
        "submitted_attempts": (
            scalar(
                conn,
                "SELECT count(*) FROM submit_attempts WHERE ts<=? AND submitted=1",
                (CUTOFF,),
            ),
            86,
        ),
        "generation": (
            scalar(conn, "SELECT max(generation) FROM alphas WHERE ts<=?", (CUTOFF,)),
            136,
        ),
        "analysis_full_pass": (
            scalar(
                conn,
                "SELECT count(*) FROM alphas WHERE ts<=? AND sharpe IS NOT NULL "
                "AND fail_count=0 AND error_count=0",
                (CUTOFF,),
            ),
            442,
        ),
        "unknown_operator": (
            scalar(
                conn,
                "SELECT count(*) FROM alphas WHERE ts<=? "
                "AND lower(error_text) LIKE '%unknown operator%'",
                (CUTOFF,),
            ),
            135,
        ),
        "unknown_variable": (
            scalar(
                conn,
                "SELECT count(*) FROM alphas WHERE ts<=? "
                "AND lower(error_text) LIKE '%unknown variable%'",
                (CUTOFF,),
            ),
            171,
        ),
        "incompatible_unit": (
            scalar(
                conn,
                "SELECT count(*) FROM alphas WHERE ts<=? "
                "AND lower(error_text) LIKE '%incompatible unit%'",
                (CUTOFF,),
            ),
            353,
        ),
    }

    for origin, expected in {
        "sweep": (1_626, 366, 112, 7),
        "crossover": (2_189, 380, 87, 8),
        "mutate": (6_465, 556, 117, 8),
        "improve": (264, 97, 34, 5),
        "ht_rescue": (172, 96, 27, 3),
        "random": (2_582, 9, 12, 6),
    }.items():
        actual = conn.execute(
            "SELECT count(*), sum(sharpe>=1.58), "
            "sum(fail_count=0 AND error_count=0), sum(submitted) "
            "FROM alphas WHERE ts<=? AND sharpe IS NOT NULL AND origin=?",
            (CUTOFF, origin),
        ).fetchone()
        checks[f"origin_{origin}"] = (tuple(actual), expected)

    round_id = scalar(
        conn,
        "SELECT id FROM rounds WHERE user_id=2 AND round_num=198 "
        "AND started_at<=? ORDER BY started_at DESC LIMIT 1",
        (CUTOFF,),
    )
    r198 = {
        row[0]: (row[1], row[2], row[3], json.loads(row[4]))
        for row in conn.execute(
            "SELECT idx, sharpe, submitted, fail_count, fail_items "
            "FROM alphas WHERE round_id=?",
            (round_id,),
        )
    }
    checks["r198_alpha2"] = (r198[2], (2.02, 0, 2, ["LOW_FITNESS", "PROD_CORRELATION"]))
    checks["r198_alpha6"] = (r198[6], (1.44, 1, 0, []))

    seoul = ZoneInfo("Asia/Seoul")
    week_start = dt.datetime(2026, 7, 28, tzinfo=seoul).timestamp()
    week_end = dt.datetime(2026, 8, 4, tzinfo=seoul).timestamp()
    week_rows = conn.execute(
        "SELECT code FROM alphas WHERE ts>=? AND ts<? AND submitted=1",
        (week_start, week_end),
    ).fetchall()
    week_rsk70 = sum("rsk70" in (code or "").lower() for (code,) in week_rows)
    checks["week_submitted"] = (len(week_rows), 17)
    checks["week_rsk70"] = (week_rsk70, 11)

    for name, (actual, expected) in checks.items():
        if actual != expected:
            raise AssertionError(f"{name}: expected {expected!r}, got {actual!r}")

    deck = ppt_text(ROOT / "docs/머신발표/GenomicWQB_머신발표.pptx")
    report = docx_text(ROOT / "docs/유전알고리즘_알파리서치_리포트.docx")
    script = (ROOT / "docs/머신발표/발표대본.md").read_text()
    combined = "\n".join((deck, report, script))

    for expected_text in [
        "35,701",
        "923",
        "submitted=1",
        "17건 중 11건",
        "unknown operator 135건",
        "unknown variable 171건",
        "unit incompatibility 353건",
    ]:
        if expected_text not in combined:
            raise AssertionError(f"required text missing: {expected_text}")

    for banned in [
        "최다 제출자",
        "3주 연속",
        "15 12 10",
        "동시 4슬롯",
        "제출 21건 중 11건",
        "문법 오류로 중단된 시뮬은 없었습니다",
        "—",
    ]:
        if banned in combined:
            raise AssertionError(f"banned or stale claim remains: {banned}")

    backend = (ROOT / "server/wqb_backend.py").read_text()
    if "_CONCURRENCY_DEFAULT = 8" not in backend:
        raise AssertionError("current backend concurrency default is not 8")

    print(f"PASS: {len(checks)} ledger/code checks and text consistency checks")


if __name__ == "__main__":
    main()
