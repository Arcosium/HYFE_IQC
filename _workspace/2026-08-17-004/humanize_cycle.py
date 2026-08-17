from pathlib import Path

from docx import Document
from pptx import Presentation


ROOT = Path("/home/arcosium/projects/GenomicWQB")
PPTX = ROOT / "docs/머신발표/GenomicWQB_머신발표.pptx"
REPORT = ROOT / "docs/유전알고리즘_알파리서치_리포트.docx"


def replace_paragraph_text(paragraph, old: str, new: str) -> bool:
    if old not in paragraph.text:
        return False
    target = paragraph.text.replace(old, new)
    if paragraph.runs:
        paragraph.runs[0].text = target
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = target
    return True


def replace_ppt(old: str, new: str) -> None:
    presentation = Presentation(PPTX)
    hits = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                hits += replace_paragraph_text(paragraph, old, new)
    if hits != 1:
        raise RuntimeError(f"PPT replacement count for {old!r}: {hits}")
    presentation.save(PPTX)


def replace_report(replacements: list[tuple[str, str]]) -> None:
    document = Document(REPORT)
    for old, new in replacements:
        hits = 0
        for paragraph in document.paragraphs:
            hits += replace_paragraph_text(paragraph, old, new)
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        hits += replace_paragraph_text(paragraph, old, new)
        if hits != 1:
            raise RuntimeError(f"report replacement count for {old!r}: {hits}")
    document.save(REPORT)


def main() -> None:
    for old, new in [
        (
            '따라서 "많이 찾는 것"과 "통과할 것을 찾는 것"은 다른 문제입니다.',
            "탐색량을 늘리는 일과 통과할 후보를 찾는 일은 다른 문제입니다.",
        ),
        (
            '핵심은 "좋은 알파를 설계하는 법"이 아니라 "제한된 슬롯을 어디에 던질지 결정하는 법"입니다. 이 발표는 그 결정 규칙을 다룹니다.',
            "핵심은 좋은 알파 설계가 아니라 제한된 슬롯의 배분입니다. 이 발표는 그 결정 규칙을 다룹니다.",
        ),
        (
            '부트캠프 기간에 검증 가능한 산출물은 IS 성과가 아니라 "몇 건을 등재시켰는가"입니다. 목적함수를 그쪽에 맞췄습니다.',
            "부트캠프 기간에 검증 가능한 산출물은 IS 성과보다 등재 건수입니다. 목적함수도 여기에 맞췄습니다.",
        ),
        (
            'LLM이 낸 수식은 변이 없이 먼저 성적을 잽니다. 그래야 "그 아이디어가 먹혔는지"를 정직하게 말할 수 있습니다. 이후 GA 재료가 됩니다.',
            "LLM이 낸 수식은 변이 없이 먼저 성적을 잽니다. 그래야 원래 아이디어의 성적을 정직하게 말할 수 있습니다. 이후 GA 재료가 됩니다.",
        ),
    ]:
        replace_ppt(old, new)

    replace_report(
        [
            (
                '질문은 "어떤 수식이 좋은가"가 아니라 "다음 여덟 슬롯을 어디에 던질 것인가"이다.',
                "연구 질문은 좋은 수식을 고르는 법보다 다음 여덟 슬롯을 어디에 배분할지에 가깝다.",
            ),
            (
                "따라서 표준 창 제약은 주된 탐색 경로에 적용된 규칙이지 전체 원장을 완전히 제한한 절대 규칙은 아니다.",
                "표준 창 제약은 주된 탐색 경로에 적용된 규칙이지 전체 원장을 완전히 제한한 절대 규칙은 아니다.",
            ),
            (
                "따라서 변경 유전자 수별 자기상관 차이를 부모와 자식이 서로 닮았다는 증거로 해석할 수 없다.",
                "변경 유전자 수별 자기상관 차이를 부모와 자식이 서로 닮았다는 증거로 해석할 수는 없다.",
            ),
            (
                "따라서 현재 자료로는 다목적 선택이 엘리트 다양성을 얼마나 높였는지 인과적으로 확인하기 어렵다.",
                "현재 자료로는 다목적 선택이 엘리트 다양성을 얼마나 높였는지 인과적으로 확인하기 어렵다.",
            ),
        ]
    )


if __name__ == "__main__":
    main()
