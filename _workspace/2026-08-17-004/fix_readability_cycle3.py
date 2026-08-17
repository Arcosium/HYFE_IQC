from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt


PPTX = Path(
    "/home/arcosium/projects/GenomicWQB/docs/머신발표/GenomicWQB_머신발표.pptx"
)


def shape(slide, name: str):
    return next(item for item in slide.shapes if item.name == name)


def set_paragraphs(shape_, texts: list[str], size: float) -> None:
    frame = shape_.text_frame
    paragraphs = list(frame.paragraphs)
    while len(paragraphs) < len(texts):
        paragraphs.append(frame.add_paragraph())
    for paragraph, text in zip(paragraphs, texts):
        if paragraph.runs:
            paragraph.runs[0].text = text
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = text
        for run in paragraph.runs:
            run.font.size = Pt(size)
    for paragraph in paragraphs[len(texts) :]:
        frame._element.remove(paragraph._element)
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.word_wrap = True


def main() -> None:
    presentation = Presentation(PPTX)

    for slide in list(presentation.slides)[1:19]:
        for name in ("TextBox 2", "TextBox 3"):
            item = shape(slide, name)
            item.height = Inches(0.45)
            item.text_frame.auto_size = MSO_AUTO_SIZE.NONE

    slide11 = presentation.slides[10]
    mapping = shape(slide11, "TextBox 12")
    mapping.height = Inches(1.25)
    for paragraph in mapping.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(10.5)

    slide19 = presentation.slides[18]
    set_paragraphs(
        shape(slide19, "TextBox 5"),
        [
            "이번 발표의 결론은 머신이 잘 돈다는 게 아니라 무엇이 실제로 효과였는지를 원장에 남겼다는 것입니다. 그 원장이 다음 설계의 근거입니다."
        ],
        12.0,
    )

    presentation.save(PPTX)


if __name__ == "__main__":
    main()
