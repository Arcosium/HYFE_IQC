from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


PPTX = Path(
    "/home/arcosium/projects/GenomicWQB/docs/머신발표/GenomicWQB_머신발표.pptx"
)
FOOTER_TOP = 6_300_000


def set_paragraph_text(paragraph, text: str) -> None:
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = text


def set_shape_paragraphs(shape, lines: list[str]) -> None:
    text_frame = shape.text_frame
    paragraphs = list(text_frame.paragraphs)
    while len(paragraphs) < len(lines):
        paragraphs.append(text_frame.add_paragraph())
    for paragraph, text in zip(paragraphs, lines):
        set_paragraph_text(paragraph, text)
    for paragraph in paragraphs[len(lines):]:
        text_frame._element.remove(paragraph._element)


def shape(slide, name: str):
    return next(item for item in slide.shapes if item.name == name)


def replace_all_text(presentation: Presentation, old: str, new: str) -> None:
    hits = 0
    for slide in presentation.slides:
        for item in slide.shapes:
            if not hasattr(item, "text_frame"):
                continue
            for paragraph in item.text_frame.paragraphs:
                if old in paragraph.text:
                    set_paragraph_text(paragraph, paragraph.text.replace(old, new))
                    hits += 1
    if hits == 0:
        raise RuntimeError(f"Text not found: {old}")


def enlarge_fonts(presentation: Presentation) -> None:
    for slide_index, slide in enumerate(presentation.slides):
        for item in slide.shapes:
            if not hasattr(item, "text_frame") or item.top >= FOOTER_TOP:
                continue
            for paragraph in item.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip() or run.font.size is None:
                        continue
                    current = run.font.size.pt
                    target = current
                    if slide_index > 0 and item.name == "TextBox 4":
                        target = max(current, 23.0)
                    elif slide_index > 0 and item.name == "TextBox 3":
                        target = max(current, 12.5)
                    elif slide_index > 0 and item.name == "TextBox 5":
                        target = max(current, 12.0)
                    elif current < 8.5:
                        target = 10.5
                    elif current < 9.5:
                        target = 11.0
                    elif current < 10.5:
                        target = 11.5
                    elif current < 11.5:
                        target = 12.0
                    if target != current:
                        run.font.size = Pt(target)


def main() -> None:
    presentation = Presentation(PPTX)

    replacements = [
        (
            "알파 리서치는 수식 문제가 아니라, 평가 예산 배분 문제입니다",
            "알파 리서치는 수식 문제가 아니라 평가 예산 배분 문제입니다",
        ),
        (
            "시뮬 1건은 대기열 포함 2~20분입니다. 동시 슬롯은 4개.",
            "시뮬 1건은 대기열 포함 2~20분입니다. 동시 슬롯은 8개.",
        ),
        (
            "하루에 물리적으로 던질 수 있는 실험은 400건 안팎입니다.",
            "대기열과 재시도를 포함한 운영 실측은 하루 평균 364건입니다.",
        ),
        (
            "동시 4슬롯을 24시간 채운 결과",
            "대기열·재시도를 포함한 운영 실측",
        ),
        ("BRAIN REST · 동시 4슬롯", "BRAIN REST · 동시 8슬롯"),
        (
            "슬롯이 4개뿐이므로, 떨어질 것이 확실한 후보를 시뮬에 태우는 것은 그 자체로 탐색 손실입니다.",
            "동시 슬롯이 8개여도 확실히 실패할 후보를 시뮬에 태우면 다음 실험이 그만큼 밀립니다.",
        ),
        (
            "문법이 보장된 고정 유전자 공간에서 렌더링합니다. 표현 자유도를 일부 포기하는 대신 문법 오류와 과도한 복잡도를 줄였습니다.",
            "고정 유전자 공간에서 렌더링하고 사전 검문으로 오류를 줄입니다. 다만 플랫폼 단계의 오류를 모두 막지는 못했습니다.",
        ),
        (
            "2·3·5·10·20·60·120·252 (8종 고정)",
            "주요 창 2·3·5·10·20·60·120·252 · 일부 경로 40",
        ),
        (
            "렌더러와 사전 검문이 문법·단위·중복을 확인합니다. 3만 5천 건을 돌리는 동안 문법 오류로 중단된 시뮬은 없었습니다. 유전자 목록이 곧 이 머신의 탐색 공간입니다.",
            "사전 검문은 오류를 줄이지만 완전하지 않습니다. 원장에는 unknown operator 135건, unknown variable 171건, unit incompatibility 353건이 남았습니다.",
        ),
        (
            "2026년 7월 28일~8월 3일 제출 21건 중 11건",
            "2026년 7월 28일~8월 3일 원장상 제출 17건 중 11건",
        ),
        (
            "제출 건수가 아니라 \"서로 prod 상관 0.7 미만인 독립 알파를 몇 개 만들었는가\"로 머신의 성적을 매깁니다.",
            "제출 건수가 아니라 프로덕션 상관 관문을 통과한 독립 계보 수로 머신의 성적을 매깁니다.",
        ),
        (
            "원장 기준 제출 성공 86건",
            "원장상 제출 성공 86건",
        ),
        (
            "submit_attempts에 제출 성공 86건이 기록됐습니다.",
            "submit_attempts에 submitted=1인 기록 86건이 남았습니다.",
        ),
        (
            "이 머신의 산출물은 제출된 알파 86건에 그치지 않습니다.",
            "이 머신의 산출물은 원장상 제출 성공 86건에 그치지 않습니다.",
        ),
    ]
    for old, new in replacements:
        replace_all_text(presentation, old, new)

    slide4 = presentation.slides[3]
    set_shape_paragraphs(
        shape(slide4, "TextBox 13"),
        [
            "① 3개월 안에는 IS Sharpe를 OS 성과로 검증하기 어렵습니다.",
            "② 제출 통과는 플랫폼이 즉시 돌려주는 반복 가능한 판정입니다.",
            "③ Power Pool·피라미드는 서로 다른 알파가 쌓여야 커집니다.",
            "④ 목적함수는 제출 통과 건수로 두고 품질은 관문 충족 여부로 봅니다.",
        ],
    )
    set_shape_paragraphs(
        shape(slide4, "TextBox 22"),
        [
            "같은 계보를 반복하면 프로덕션 상관도 함께 오릅니다.",
            "8월 13~14일 태그 134건 중 53건에 값이 남았습니다. 중앙값은 0.9028, 최소값은 0.7111이었습니다.",
            "다음 목표는 제출량보다 독립 계보 수입니다. 19쪽에서 다룹니다.",
        ],
    )

    slide9 = presentation.slides[8]
    set_shape_paragraphs(
        shape(slide9, "TextBox 11"),
        [
            "FAIL인데도 제출됐다면 그 체크는 단독 차단 근거가 아닙니다.",
            "403 본문에는 한 알파의 여러 FAIL이 함께 실립니다.",
            "성공 이력이 거절 목록보다 강한 증거입니다.",
        ],
    )
    set_shape_paragraphs(
        shape(slide9, "TextBox 14"),
        [
            "최근순 비교만 쓰면 거절이 많은 체크가 모두 하드 규칙으로 굳습니다.",
            "실측에서는 거절 135건·성공 23건인데 soft=[]였고 LOW_FITNESS가 하드로 분류됐습니다.",
            "해당 부류는 fitness 0.26~0.86인 22건이 모두 제출 성공한 기록이었습니다.",
        ],
    )

    slide10 = presentation.slides[9]
    set_shape_paragraphs(
        shape(slide10, "TextBox 12"),
        [
            "동시에 바꾼 유전자가 늘수록 부모 대비 개선율이 낮아졌습니다.",
            "1개 변경 20.3% → 6개 이상 변경 6.2%",
            "표본은 유효 부모·자식 쌍 10,119건입니다.",
            "인과효과가 아니라 원인을 한 축에 귀속하기 쉬운 설계라는 뜻입니다.",
        ],
    )

    slide11 = presentation.slides[10]
    set_shape_paragraphs(
        shape(slide11, "TextBox 15"),
        [
            "고회전 관문 미달은 FAIL이 아니라 WARNING이라 fail_items에 잡히지 않았습니다.",
            "지시 선택에 부모의 실측 지표를 더해 숨은 병목을 보완했습니다.",
        ],
    )
    set_shape_paragraphs(
        shape(slide11, "TextBox 17"),
        [
            "누적 관측에서 boost 36.0%, smooth 0.5%로 약 72배 차이가 났습니다.",
            "통제 실험은 아니며 성공률 표가 다음 변이의 우선순위를 바꿉니다.",
        ],
    )

    slide12 = presentation.slides[11]
    set_shape_paragraphs(
        shape(slide12, "TextBox 18"),
        [
            "비지배 정렬로 프런트를 나누고 같은 프런트에서는 혼잡 거리가 큰 개체를 남깁니다.",
            "성능이 조금 낮아도 목적 공간의 다른 영역을 지키는 엘리트를 보존합니다.",
        ],
    )
    set_shape_paragraphs(
        shape(slide12, "TextBox 23"),
        [
            "전 기간에서 엘리트를 뽑으면 한 번 잘 나온 가계가 자리를 독점합니다.",
            "최근 N일 창으로 오래된 승자를 밀어내지만 세대 깊이가 얕아지는 비용이 있습니다.",
        ],
    )
    set_shape_paragraphs(
        shape(slide12, "TextBox 29"),
        [
            "family 정보가 남은 제출작 25건 중 model 14건과 pv 8건이 88%를 차지했습니다.",
            "후보가 두 계열에서만 나오면 다목적 선택도 그 범위 안에서만 작동합니다.",
            "다음 비교에서는 생성 팔레트와 엘리트 풀 분포를 함께 기록합니다.",
        ],
    )

    slide14 = presentation.slides[13]
    set_shape_paragraphs(
        shape(slide14, "TextBox 18"),
        [
            "프로덕션 상관 조회는 167초를 기다려도 빈 본문이었습니다.",
            "당시 실값은 제출 시도의 403 본문에서 가장 안정적으로 확인됐습니다.",
        ],
    )
    set_shape_paragraphs(
        shape(slide14, "TextBox 23"),
        [
            "형제 알파가 OS에 오르면 같은 계보의 prod 상관이 수 시간 안에 0.75+로 나타났습니다.",
            "전날 판정을 재사용하지 않고 당일 기준을 다시 받습니다.",
        ],
    )
    set_shape_paragraphs(
        shape(slide14, "TextBox 26"),
        [
            "쿼터가 찬 뒤의 통과작은 다음 리셋 직후 발사하도록 대기 큐에 넣습니다.",
            "이 안전망이 없던 8월 12일에는 통과작 27건이 큐에 남지 않았습니다.",
            "다음 날 복원한 후보 중 한 건이 제출됐습니다.",
        ],
    )

    slide15 = presentation.slides[14]
    set_shape_paragraphs(
        shape(slide15, "TextBox 12"),
        [
            "무작위 탐색은 2,582건 중 9건이 S≥1.58에 도달했습니다. 수율은 0.3%입니다.",
            "교차·스윕·개선·구제 경로의 관측 수율은 17.4~55.8%였습니다.",
            "무작위 탐색은 새 데이터셋에 처음 들어가는 기준선으로 남깁니다.",
        ],
    )

    slide16 = presentation.slides[15]
    set_shape_paragraphs(
        shape(slide16, "TextBox 12"),
        [
            "① LOW_FITNESS·LOW_SHARPE는 약한 신호에서 생기는 상시 배경 잡음입니다.",
            "② LOW_2Y·IS_LADDER·SUB_UNIVERSE와 리전별 Sharpe도 반복됐습니다.",
            "③ PROD_CORRELATION은 제출 직전에 주로 측정돼 누적 건수가 작습니다.",
        ],
    )
    set_shape_paragraphs(
        shape(slide16, "TextBox 15"),
        [
            "누적 빈도와 현재 병목은 다릅니다.",
            "8월 13~14일 태그 134건 중 prod 상관값은 53건에 남았고 중앙값은 0.9028이었습니다.",
            "0.7 미만은 없었습니다. 이 시점의 병목은 약한 신호와 높은 유사도가 함께 만든 결과입니다.",
        ],
    )

    slide17 = presentation.slides[16]
    set_shape_paragraphs(
        shape(slide17, "TextBox 14"),
        [
            "중립화만 INDUSTRY → STATISTICAL로 바꾸자 Sharpe가 1.72 → 2.47로 올랐습니다.",
            "감쇠만 0 → 8로 바꾸자 회전율이 100.5% → 48.5%로 낮아졌습니다.",
            "한 축만 달라 결과를 그 변화에 귀속할 수 있었습니다.",
        ],
    )
    set_shape_paragraphs(
        shape(slide17, "TextBox 20"),
        [
            "2026년 7월 28일~8월 3일 원장상 제출 17건 중 11건이 rsk70에서 나왔습니다.",
            "당시 시딩 축 26개가 모두 rsk70이어서 생성 단계부터 치우쳤습니다.",
            "선택 효과만 따로 확인할 지표는 없었습니다.",
        ],
    )
    set_shape_paragraphs(
        shape(slide17, "TextBox 26"),
        [
            "8월 12일 하루 4건 상한을 채운 뒤 통과작 27건이 큐에 남지 않았습니다.",
            "REGULAR_SUBMISSION을 알파 결함으로 오분류해 익일 큐를 우회한 것이 원인이었습니다.",
            "다음 날 복원한 후보 중 한 건(a#39827, S 1.97)이 제출됐습니다.",
        ],
    )

    slide18 = presentation.slides[17]
    set_shape_paragraphs(
        shape(slide18, "TextBox 11"),
        [
            "8월 12일 같은 계보 네 건이 OS에 등재된 뒤 그 계보의 prod 상관이 함께 올랐습니다.",
            "다음 날 태그 134건 중 53건에 값이 남았고 최소값도 0.7111이었습니다.",
            "같은 골격을 반복하면 스스로 제출 가능성을 깎습니다.",
        ],
    )
    set_shape_paragraphs(
        shape(slide18, "TextBox 14"),
        [
            "피라미드는 리전·지연·데이터셋 카테고리 조합에 알파 세 건을 제출하면 하나로 집계됩니다.",
            "7월 28일~8월 3일 원장상 제출 17건 중 11건이 rsk70에 몰렸습니다.",
            "미달 칸을 우선하고 같은 데이터셋 안에서 짝을 고르도록 생성 정책을 바꿨습니다.",
        ],
    )

    slide19 = presentation.slides[18]
    set_shape_paragraphs(
        shape(slide19, "TextBox 14"),
        [
            "같은 스냅샷에서 Analyst 4,218필드와 Fundamental 4,428필드에 새 계보를 심습니다.",
            "제출작의 88%가 두 패밀리에 몰린 쏠림을 생성 단계에서 줄입니다.",
        ],
    )

    enlarge_fonts(presentation)
    presentation.save(PPTX)


if __name__ == "__main__":
    main()
