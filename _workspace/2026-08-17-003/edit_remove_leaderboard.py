from pathlib import Path

from pptx import Presentation


PPTX = Path(
    "/home/arcosium/projects/GenomicWQB/docs/머신발표/GenomicWQB_머신발표.pptx"
)


def set_paragraph_text(paragraph, text: str) -> None:
    """Replace text while retaining the first run's existing formatting."""
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = text


def remove_paragraph(text_frame, paragraph) -> None:
    text_frame._element.remove(paragraph._element)


def main() -> None:
    presentation = Presentation(PPTX)

    cover = presentation.slides[0]
    cover_stats = next(shape for shape in cover.shapes if shape.name == "TextBox 7")
    for paragraph in list(cover_stats.text_frame.paragraphs):
        if "주간 최다 제출자" in paragraph.text:
            remove_paragraph(cover_stats.text_frame, paragraph)

    problem = presentation.slides[3]
    result_label = next(shape for shape in problem.shapes if shape.name == "TextBox 16")
    result_title = next(shape for shape in problem.shapes if shape.name == "TextBox 17")
    result_body = next(shape for shape in problem.shapes if shape.name == "TextBox 18")

    set_paragraph_text(result_label.text_frame.paragraphs[0], "운영 결과")
    set_paragraph_text(result_title.text_frame.paragraphs[0], "원장 기준 제출 성공 86건")

    replacement = [
        "2026년 5월 9일부터 8월 14일 16시 20분까지 submit_attempts에 제출 성공 86건이 기록됐습니다.",
        "목적함수와 결과는 같은 기준시각의 운영 원장으로 검증했습니다.",
        "주차별 비교 대신 누적 제출 성공과 통과 조건을 봅니다.",
    ]
    paragraphs = result_body.text_frame.paragraphs
    if len(paragraphs) != len(replacement):
        raise RuntimeError(f"Unexpected result-body paragraph count: {len(paragraphs)}")
    for paragraph, text in zip(paragraphs, replacement):
        set_paragraph_text(paragraph, text)

    presentation.save(PPTX)


if __name__ == "__main__":
    main()
