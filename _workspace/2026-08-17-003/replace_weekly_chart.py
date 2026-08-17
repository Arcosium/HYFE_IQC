from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt


PPTX = Path(
    "/home/arcosium/projects/GenomicWQB/docs/머신발표/GenomicWQB_머신발표.pptx"
)


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = 0
    paragraph.space_after = 0
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def main() -> None:
    presentation = Presentation(PPTX)
    slide = presentation.slides[3]

    chart = next(
        shape
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        and shape.left == 566928
        and shape.top == 3744837
    )
    x, y, w, h = chart.left, chart.top, chart.width, chart.height
    remove_shape(chart)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(247, 249, 252)
    panel.line.color.rgb = RGBColor(207, 216, 228)
    panel.line.width = Pt(0.8)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, 41148, h)
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(11, 79, 158)
    accent.line.fill.background()

    add_text(
        slide,
        x + 219456,
        y + 118872,
        w - 438912,
        228600,
        "같은 기준시각의 원장 스냅샷",
        10.5,
        (11, 79, 158),
        bold=True,
    )
    add_text(
        slide,
        x + 219456,
        y + 329184,
        w - 438912,
        210312,
        "2026.05.09~08.14 16:20 KST",
        8.5,
        (90, 90, 90),
    )

    metrics = [
        ("98일", "운영 기간"),
        ("35,701회", "시뮬레이션"),
        ("923회", "라운드"),
        ("86건", "제출 성공"),
    ]
    inner_x = x + 219456
    gap = 109728
    card_y = y + 649224
    card_h = 1188720
    card_w = (w - 438912 - gap * 3) // 4
    for index, (value, label) in enumerate(metrics):
        card_x = inner_x + index * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_x, card_y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(224, 230, 238)
        card.line.width = Pt(0.7)
        add_text(
            slide,
            card_x + 45720,
            card_y + 201168,
            card_w - 91440,
            420624,
            value,
            19,
            (26, 26, 26),
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            card_x + 45720,
            card_y + 685800,
            card_w - 91440,
            283464,
            label,
            9,
            (90, 90, 90),
            align=PP_ALIGN.CENTER,
        )

    add_text(
        slide,
        x + 219456,
        y + h - 374904,
        w - 438912,
        210312,
        "alphas · rounds · submit_attempts 집계",
        8,
        (112, 120, 132),
    )

    presentation.save(PPTX)


if __name__ == "__main__":
    main()
