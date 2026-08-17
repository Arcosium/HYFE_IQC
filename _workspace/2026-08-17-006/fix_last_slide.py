from pathlib import Path

from pptx import Presentation


PPTX = Path(
    "/home/arcosium/projects/GenomicWQB/docs/머신발표/GenomicWQB_머신발표.pptx"
)


def set_paragraph_text(paragraph, text: str) -> None:
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = text


def main() -> None:
    presentation = Presentation(PPTX)
    last_slide = presentation.slides[-1]
    body = next(shape for shape in last_slide.shapes if shape.name == "TextBox 4")
    paragraphs = list(body.text_frame.paragraphs)
    if len(paragraphs) != 3:
        raise RuntimeError(f"unexpected paragraph count: {len(paragraphs)}")

    set_paragraph_text(
        paragraphs[0],
        "이 머신의 산출물은 원장상 제출 성공 86건에 그치지 않습니다.",
    )
    set_paragraph_text(
        paragraphs[1],
        "어떤 유전자가 어느 환경에서 작동했고 무엇이 실패했는지를 남긴 35,701건의",
    )
    set_paragraph_text(
        paragraphs[2],
        "실험 원장이 같은 무게의 결과물입니다.",
    )

    presentation.save(PPTX)


if __name__ == "__main__":
    main()
