from pathlib import Path

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import RGBColor
from pptx import Presentation
from pptx.util import Inches


ROOT = Path("/home/arcosium/projects/GenomicWQB")
PPTX = ROOT / "docs/머신발표/GenomicWQB_머신발표.pptx"
REPORT = ROOT / "docs/유전알고리즘_알파리서치_리포트.docx"


# Handmade table regions in inches, keyed by one-based slide number.
TABLE_REGIONS = {
    5: (5.20, 6.90),
    6: (2.10, 5.10),
    9: (4.90, 6.90),
    10: (5.50, 6.90),
    13: (2.10, 4.70),
    15: (5.00, 6.90),
    18: (5.00, 6.90),
}


def flatten_presentation_tables() -> int:
    presentation = Presentation(PPTX)
    changed = 0

    for slide_number, (top_min, top_max) in TABLE_REGIONS.items():
        slide = presentation.slides[slide_number - 1]
        for shape in slide.shapes:
            if not shape.name.startswith("Rectangle"):
                continue
            top = shape.top / 914_400
            height = shape.height / 914_400
            if not (top_min <= top <= top_max and height <= 0.60):
                continue
            for effect_ref in shape._element.xpath("./p:style/a:effectRef"):
                if effect_ref.get("idx") != "0":
                    effect_ref.set("idx", "0")
                    changed += 1

    # Keep a clear gap between the closing rule and the thank-you line.
    last_slide = presentation.slides[-1]
    thanks = next(shape for shape in last_slide.shapes if shape.name == "TextBox 6")
    thanks.top = Inches(5.50)
    thanks.height = Inches(0.45)

    presentation.save(PPTX)
    return changed


def set_xml_runs_black(root) -> int:
    count = 0
    for run in root.xpath(".//w:r"):
        rpr = run.find(qn("w:rPr"))
        if rpr is None:
            rpr = OxmlElement("w:rPr")
            run.insert(0, rpr)
        color = rpr.find(qn("w:color"))
        if color is None:
            color = OxmlElement("w:color")
            rpr.append(color)
        color.set(qn("w:val"), "000000")
        for attr in ("themeColor", "themeTint", "themeShade"):
            color.attrib.pop(qn(f"w:{attr}"), None)
        count += 1
    return count


def format_report() -> tuple[int, int]:
    document = Document(REPORT)

    for style in document.styles:
        if hasattr(style, "font"):
            style.font.color.rgb = RGBColor(0, 0, 0)

    centered_cells = 0
    for table in document.tables:
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        for row in table.rows:
            for cell in row.cells:
                cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
                for paragraph in cell.paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                centered_cells += 1

    black_runs = set_xml_runs_black(document._element)
    seen_parts = set()
    for section in document.sections:
        for story in (section.header, section.footer):
            part_name = str(story.part.partname)
            if part_name in seen_parts:
                continue
            seen_parts.add(part_name)
            black_runs += set_xml_runs_black(story._element)

    document.save(REPORT)
    return black_runs, centered_cells


def main() -> None:
    flattened = flatten_presentation_tables()
    black_runs, centered_cells = format_report()
    print(
        f"flattened_table_shapes={flattened} "
        f"black_runs={black_runs} centered_cells={centered_cells}"
    )


if __name__ == "__main__":
    main()
