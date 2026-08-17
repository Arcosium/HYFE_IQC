from __future__ import annotations

import json
import re
import sqlite3
from collections import Counter
from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib import font_manager
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Mm, Pt, RGBColor
from pptx import Presentation


ROOT = Path('/home/arcosium/projects/GenomicWQB')
DOCX = ROOT / 'docs/유전알고리즘_알파리서치_리포트.docx'
PPTX = ROOT / 'docs/머신발표/GenomicWQB_머신발표.pptx'
FAIL_FIG = ROOT / 'docs/머신발표/figs/fig5_failures.png'
DB = ROOT / 'data/hyfe_iqc.db'
CUTOFF = 1786692000  # 2026-08-14 16:20:00 KST
ORIGINS = ('sweep', 'crossover', 'mutate', 'improve', 'ht_rescue', 'random')


def report_edits() -> None:
    doc = Document(DOCX)

    replacements = {
        1: '한 축 변이와 다목적 선택의 탐색 수율·상관 구조 분석',
        2: 'Genetic Algorithms for Alpha Research: Search Yield and Correlation Structure under Single-Axis Mutation and Multi-Objective Selection',
        6: '알파 리서치는 수식을 잘 쓰는 문제가 아니라 평가 예산을 어디에 던질지 정하는 문제다. 후보의 품질은 WorldQuant BRAIN 시뮬레이션을 돌린 뒤에야 드러나고 동시 실행 슬롯은 네 개뿐이라 전수 탐색은 성립하지 않는다. 본 연구는 알파를 고정 유전자 공간의 유전체로 표현하고 선택·교차·변이를 반복하는 자동 발굴 시스템 GenomicWQB를 98일간 운영한 원장을 분석한다. 시뮬레이션 35,701건, 923라운드, 고유 표현식 26,404종이 분석 대상이다.',
        7: '두 가지를 물었다. 첫째, 여러 유전자를 한꺼번에 바꾸는 방식과 비교할 때 한 축 스윕과 실패 원인별 정향변이의 개선율은 어떻게 달라지는가. 둘째, 다목적 선택을 적용한 뒤 엘리트와 제출작의 다양성은 어떤 모습으로 나타나는가.',
        8: '동시에 바꾼 유전자가 늘수록 부모 대비 Sharpe 개선율은 낮아졌다(1개 20.3%, 6개 이상 6.2%, 유효 부모·자식 쌍 n=10,119). Sharpe가 기록된 주요 자동 생성 경로 여섯 개에서는 한 축 스윕의 S≥1.58 도달률이 22.5%, 무작위 탐색은 0.3%였다. 자기상관은 후보와 기존 제출 풀의 유사도를 나타낸다. 한 축 변경군의 중앙값은 0.3456, 여섯 축 이상 변경군은 0.2771이었다. 이는 부모와 자식의 직접 상관이 아니라, 변경 폭과 기존 제출 풀 유사도 사이의 관측 관계다.',
        9: '다목적 선택은 목적 공간의 여러 영역을 보존하도록 설계됐다. 다만 현재 원장만으로 적용 전후의 다양성 변화를 인과적으로 분리할 수는 없다. family 정보가 남은 제출작 25건 가운데 22건(88%)이 model과 pv 두 패밀리에 몰렸다. 선택 규칙과 함께 생성 팔레트도 바뀌어 왔으므로, 확인 가능한 결론은 제출작의 쏠림이 여전히 남았다는 점이다. 다양성은 선택뿐 아니라 생성 단계에서도 확보해야 한다.',
        10: '본 시스템은 알파의 사후 수익률이 아니라 제출 통과 건수를 운영 목표로 삼았고, 부트캠프 주간 리더보드에서 최다 제출자 상을 3주 연속 받았다. 그 선택이 남긴 비용도 함께 기록한다. 8월 13~14일 원장의 `[13시전 수동]`·`[13시전 블라스트]` 태그 134건 가운데 프로덕션 상관값이 기록된 53건의 중앙값은 0.9028, 최소값은 0.7111이었다. 컷 0.7 아래는 한 건도 없었다.',
        13: 'JEL 분류: C61, C63, G11, G17',
        16: '알파 수식 하나는 데이터필드, 시계열 변환, 결합 방식, 룩백, 중립화, 감쇠, 이상치 처리 같은 여러 선택의 조합이다. 선택지가 하나 늘 때마다 조합은 곱으로 늘어난다. 2026년 8월 14일 코드와 USA/D1/TOP3000 Matrix 카탈로그를 기준으로 데이터필드 4,593종, 변환 6종, 결합 7종, 중립화 11종, 표준 룩백 8종이 있다.',
        18: '그래서 이 연구는 알파 리서치를 수식 설계 문제가 아니라 예산 배분 문제로 본다. 질문은 "어떤 수식이 좋은가"가 아니라 "다음 네 슬롯을 어디에 던질 것인가"이다. 유전 알고리즘은 넓은 탐색 공간을 재현 가능한 실험 단위로 나누고, 성과가 좋은 구조를 다음 세대로 넘기며, 실패에서도 다음 실험 방향을 얻는 절차를 제공한다[2,3].',
        19: '분석 대상인 GenomicWQB는 2026년 5월 9일부터 라운드를 운영해 왔다. LLM은 가설과 전략 스펙을 만들고, 유전 알고리즘은 반복 탐색과 계보 관리를 맡으며, 밴딧은 슬롯을 배분한다. 평가는 WorldQuant BRAIN 시뮬레이션으로 받는다. 2026년 8월 14일 16시 20분(KST)까지 98일간 시뮬레이션 35,701건과 923라운드가 쌓였고 유전체·부모 ID·세대·바뀐 유전자가 원장에 남았다.',
        20: '이 시스템은 목적함수를 알파의 사후 수익률이 아니라 제출 통과 건수로 잡았다. 3개월 부트캠프 안에서는 IS Sharpe를 OS 성과로 검증하기 어렵지만, 제출 통과는 플랫폼이 그날 돌려주는 판정이라 즉시 반복 측정할 수 있다. 같은 기간 주간 최다 제출자 상을 3주 연속 받았고, 상관 포화도 함께 나타났다.',
        26: '유전 알고리즘의 기본 요소는 개체, 유전체, 적합도, 선택, 교차, 변이다[2,3]. 알파 리서치의 언어로 옮기면 개체는 알파 후보 하나, 유전체는 그 알파를 만드는 설계 변수의 묶음이다. 적합도는 시뮬레이션 지표, 선택은 다음 라운드의 엘리트를 고르는 절차, 교차는 두 엘리트의 유전자를 섞는 것, 변이는 유전자 일부를 바꾸는 것이다.',
        27: '유전 프로그래밍과는 구분된다[2]. 유전 프로그래밍은 구문 트리를 진화시키지만 GenomicWQB는 문법과 연구 관행을 반영한 고정 유전자 공간에서 수식을 렌더링한다. 표현 자유도를 일부 포기하는 대신 문법 오류와 과도한 복잡도를 줄였다. 운영 원장 35,701건에는 문법 오류로 중단된 시뮬레이션이 없다.',
        33: '룩백은 2·3·5·10·20·60·120·252의 여덟 종으로 제한했다. 임의의 창을 넓게 훑기보다 단순하고 설명 가능한 값을 쓰라는 BRAIN 권고를 따른 것이다[9]. 표현력은 줄지만 탐색 과정의 과최적화 위험을 낮출 수 있다.',
        36: '선택은 연속 적합도와 NSGA-II를 함께 쓴다[1]. 연속 적합도는 전 체크 통과 이전의 근접 실패에도 선택압을 준다. 전 체크 통과는 `fail_count=0`, Sharpe 기록, `error_count=0`을 모두 만족한 442건으로 전체의 1.2%다. NSGA-II는 Sharpe, Fitness, 제출 경로 점수, self-correlation, 최근 2년 Sharpe를 목적축으로 삼아 파레토 면과 혼잡 거리를 계산한다. Fitness에는 회전율 효율이 반영된다.',
        38: '변이는 완전한 무작위가 아니다. 부모의 실패 항목과 실측 지표를 보고 바꿀 유전자 범주를 정하고, 관측이 쌓이면 Thompson sampling으로 지시별 성공 확률을 갱신한다[4,5]. 별도의 스윕 슬롯은 중립화나 감쇠처럼 결과에 큰 영향을 주는 축을 하나씩만 바꾼다. 차이를 한 축으로 제한하면 여러 축을 함께 바꿀 때보다 결과를 해석하기 쉽다.',
        40: '5. RQ1: 변경 폭과 개선율·기존 제출 풀 유사도의 관계',
        42: '부모와 자식의 Sharpe가 모두 기록된 유효 대응 표본 10,119건을 동시에 바꾼 유전자 수로 나누고, 자식 Sharpe가 부모보다 높은 비율을 계산했다. 변경 폭이 커질수록 개선율은 낮아졌다.',
        45: '한 개만 바꾼 자식의 20.3%가 부모보다 좋아졌고 여섯 개 이상을 바꾼 자식은 6.2%였다. 이 값은 변경 폭과 개선 여부의 기술통계다. 여러 축을 함께 바꾸면 어떤 변화가 결과와 연결됐는지 분리하기 어려우므로, 한 축 스윕은 후속 실험의 해석 가능성을 높이는 데 유리하다.',
        47: 'Sharpe가 기록된 주요 자동 생성 경로 여섯 개(n=13,298)를 대상으로 표준 제출컷(S≥1.58) 도달 비율을 비교했다. 전체 원장과 달리 기원이 없거나 다른 레거시 경로의 행은 이 표에서 제외했다.',
        51: '무작위 탐색은 2,582건 가운데 아홉 건이 컷에 도달했다. 교차·스윕·개선·구제 경로는 17.4%에서 55.8%였다. 일반 변이는 부모를 쓰지만 수율은 8.6%였으므로, 부모 유무만으로 차이를 설명할 수는 없다. 무작위 탐색은 신규 데이터셋에 처음 진입하는 역할이 있어 수율이 낮아도 탐색 폭을 위해 남겼다.',
        55: '지시별 개선율은 boost 36.0%(유효 대응 표본 631건), smooth 0.5%(185건)로 약 72배 차이가 났다. 이는 통제 실험의 인과효과가 아니라 누적 운영 관측이다. 다만 (실패 범주 × 변이 지시) 성공률을 계속 갱신하면 고정 규칙보다 최근 기록을 반영할 수 있다[4,5].',
        56: '5.4 자기상관은 부모·자식 상관이 아니다',
        57: '자기상관은 후보의 PnL을 사용자가 과거에 제출한 OS 적격 알파 풀과 비교한 최대 상관값이다. 따라서 변경 유전자 수별 자기상관 차이를 부모와 자식이 서로 닮았다는 증거로 해석할 수 없다. 여기서는 변경 폭에 따라 후보가 기존 제출 풀과 얼마나 비슷했는지를 기술한다.',
        59: '한 축 변경군의 자기상관 중앙값은 0.3456, 여섯 축 이상 변경군은 0.2771이었다. 생성 경로별로는 스윕 0.3766, 무작위 탐색 0.1780이었다. 각 값은 후보와 기존 제출 풀의 최대 상관을 요약한 것이며, 집단별 표본 구성과 체크가 기록되는 조건이 달라 직접적인 인과 비교에는 쓸 수 없다.',
        60: 'RQ1에 대해서는 두 가지를 구분해 답해야 한다. 한 축 변경군은 관측 개선율이 높고 변화의 원인을 추적하기 쉽다. 동시에 기존 제출 풀과의 자기상관 중앙값도 더 높았다. 한 축 변이가 탈상관을 악화시켰다고 단정할 수는 없지만, 개선율만으로 탐색 전략을 평가해서도 안 된다.',
        61: '실무 병목은 8월 13~14일 제출 원장에서 확인된다. `[13시전 수동]` 11건과 `[13시전 블라스트]` 123건을 합친 134건 가운데 프로덕션 상관값이 기록된 것은 53건이었다. 중앙값은 0.9028, 최소값은 0.7111로 컷 0.7 아래가 없었다. 이 시점에는 약한 신호뿐 아니라 기존 OS 알파와의 높은 유사도도 제출을 막고 있었다.',
        62: '6. RQ2: 다목적 선택 뒤에도 제출작 쏠림은 남았다',
        64: '단일 가중합으로 엘리트를 고르면 특정 유형으로 수렴할 수 있다. NSGA-II는 비지배 정렬로 파레토 프런트를 나누고, 같은 프런트 안에서는 혼잡 거리가 큰 개체를 보존한다[1]. 최근성 윈도우도 더해 오래된 승자가 계속 자리를 차지하는 것을 줄였다.',
        65: '6.2 결과: 선택 효과는 분리하기 어렵고 제출작은 쏠렸다',
        68: '원장은 세대 깊이 g136, 설정 조합 2,227종, 고유 표현식 26,404종을 기록한다. 그러나 이 규모만으로 엘리트 풀의 다양성이 개선됐다고 말할 수는 없다. family 정보가 남은 제출작 25건 가운데 model 14건, pv 8건, imbalance 2건, news 1건이었다. model과 pv가 22건(88%)을 차지했다.',
        69: '생성 팔레트의 쏠림도 확인됐다. 자동 시딩 축 목록 26개가 모두 rsk70 데이터셋 소속이었고, 2026년 7월 28일부터 8월 3일까지 제출 21건 가운데 11건이 그 데이터셋에서 나왔다. 다만 같은 기간 선택 규칙과 시딩 전략이 함께 바뀌었으므로 원인을 한 단계에만 돌릴 수는 없다.',
        70: '따라서 현재 자료로는 다목적 선택이 엘리트 다양성을 얼마나 높였는지 인과적으로 확인하기 어렵다. 확인되는 사실은 제출작의 패밀리 분포가 두 계열에 치우쳤다는 점이다. 다음 비교부터는 엘리트 풀의 패밀리 엔트로피와 교체율을 버전별로 저장해야 선택 효과를 따로 평가할 수 있다.',
        82: '둘째, 시스템이 계속 개선됐다. 98일 동안 게이트 규칙, 선택 알고리즘, 시딩 전략이 여러 차례 바뀌었다. 기원별 수율과 변경 폭별 개선율은 운영 자료의 기술통계이며 통제 실험의 효과 추정치가 아니다. 같은 리전·지연·유니버스 안에서 비교했지만 시점 효과를 완전히 제거하지 못했다.',
        83: '셋째, 자기상관 표본이 전체의 일부다. 후보가 일정 수준을 넘어 플랫폼 체크를 받았을 때만 값이 기록되므로 선택 편의가 있다. 또한 자기상관은 부모·자식 상관이 아니라 후보와 기존 제출 풀의 최대 상관이다. 여섯 개 이상 변경 구간의 표본은 78건에 불과하다.',
        84: '넷째, 제출된 알파만 보면 생존자 편향이 생긴다. 제출 통과를 운영 목표로 삼았으므로 이 편향은 설계에 내재한다. IS 성과가 실거래 수익을 보장하지 않으며 반복 탐색이 IS 과적합을 키울 수 있다는 점도 남는다[8].',
        87: '두 질문의 답은 모두 제한적이다. 한 축 변경군은 개선율이 높고 결과를 해석하기 쉬웠지만 기존 제출 풀과의 자기상관도 더 높았다. 다목적 선택은 다양성을 보존하도록 설계됐으나, 현재 원장만으로 적용 효과를 분리할 수 없었고 제출작의 패밀리 쏠림도 남았다. 다음 실험에서는 생성 팔레트뿐 아니라 엘리트 풀의 분포와 교체율도 함께 기록해야 한다.',
        88: '제출 통과 건수를 운영 목표로 둔 기간에 주간 최다 제출자 상을 3주 연속 받았다. 동시에 같은 계보를 반복해 제출하면서 프로덕션 상관이 높아졌다. 탐색 시스템의 목표는 자신이 만든 결과에 따라 제약 조건을 바꾸기도 한다. 이 되먹임을 기록하지 않으면 시스템은 이미 소진된 계보에 계속 예산을 쓸 수 있다.',
        98: '[8] Bailey, D. H., Borwein, J., López de Prado, M., & Zhu, Q. J. (2014). Pseudo-Mathematics and Financial Charlatanism: The Effects of Backtest Overfitting on Out-of-Sample Performance. Notices of the AMS, 61(5), 458–471.',
        99: '[9] WorldQuant BRAIN. Recommended Practices & Consultant Documentation. 2026-08-14 인출.',
        101: '[11] 분석 원장: GenomicWQB 운영 데이터베이스(SQLite), alphas·rounds·submit_attempts·bandit_arms·strategy_specs 테이블, 2026-08-14 16:20 KST 스냅샷.',
        103: '본 보고서의 집계 시각은 2026년 8월 14일 16시 20분(KST)이며 대상 기간은 2026년 5월 9일부터 같은 날까지다. 전체 규모는 alphas 35,701행, rounds 923행, 고유 code 26,404종, 고유 settings_fp 2,227종이다. 제출 86건은 submit_attempts의 성공 기록을 센 값이다. 기원별 수율 표는 Sharpe가 기록된 주요 자동 생성 경로 여섯 개(n=13,298)만 포함한다. 전 체크 통과 442건은 Sharpe 기록, fail_count=0, error_count=0을 모두 만족한 행이다. 계정 정보와 자격증명은 추출 대상에서 제외했다.',
        104: '코드 근거는 genome_models.py(유전체와 렌더러), selection.py(NSGA-II와 연속 적합도), reward.py(밴딧 보상), mutation_learn.py(정향변이 학습), submit_push.py(자동 시딩과 상관 완화), gate_watch.py(게이트 실측 학습), worker.py(라운드 루프), db.py(원장 스키마)다.',
    }
    for idx, text in replacements.items():
        doc.paragraphs[idx].text = text

    # 표 0: 연구 질문
    doc.tables[0].cell(2, 1).text = '다목적 선택을 적용한 뒤 엘리트와 제출작의 다양성은 어떤 모습으로 나타나는가? 현재 자료에서 선택 효과와 생성 팔레트의 영향을 어디까지 구분할 수 있는지도 함께 본다.'

    # 표 1: 유전체 범위와 2026-08-14 코드 스냅샷
    t = doc.tables[1]
    t.cell(1, 2).text = 'USA/D1/TOP3000 Matrix 4,593 · 변환 6 · 결합 7'
    t.cell(3, 2).text = '중립화 11 · 감쇠 6'
    t.cell(4, 1).text = 'truncation · nan_handling · winsor_std · sign · universe · decay_style'
    t.cell(4, 2).text = '절단 3 · 결측 2 · 윈저 3 · 부호 2 · 유니버스 6 · 감쇠 방식 2'

    # 표 2: 개선율 계산에 실제 사용한 유효 부모·자식 대응 표본 수
    paired = ('3,913', '2,176', '1,421', '1,047', '737', '825')
    for row, value in zip(doc.tables[3].rows[1:], paired):
        row.cells[1].text = value

    # 부록 표의 집계 정의 보완
    ap = doc.tables[5]
    ap.cell(2, 2).text = '부모·자식 Sharpe가 모두 기록된 대응 표본에서 자식 Sharpe − 부모 Sharpe > 0 비율을 계산하고 변경 유전자 수로 층화'
    ap.cell(3, 2).text = 'Sharpe가 기록된 주요 자동 생성 경로 6개(n=13,298)에서 S≥1.58, fail_count=0·error_count=0, 제출 건수 집계'
    ap.cell(5, 1).text = 'alphas.metrics → self_correlation'
    ap.cell(5, 2).text = '후보와 기존 제출 풀의 최대 상관 중앙값과 0.7 미만 비율'
    ap.cell(6, 2).text = 'submit_attempts 성공 86건과 주별 리더보드 집계, 거절 상태 문자열 파싱'
    ap.cell(7, 1).text = 'alphas.fail_items · submit_attempts.submit_status'
    ap.cell(7, 2).text = 'FAIL 체크 빈도와 제출 단계의 rejected·prod_corr 상태를 구분해 집계'

    # 문서 전체에서 산문용 대시와 붙어 있는 문장부호를 정리한다.
    for p in doc.paragraphs:
        if '—' in p.text or re.search(r'\.(?=[가-힣])', p.text):
            p.text = re.sub(r'\.(?=[가-힣])', '. ', p.text.replace('—', ':'))
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    if '—' in p.text or re.search(r'\.(?=[가-힣])', p.text):
                        p.text = re.sub(r'\.(?=[가-힣])', '. ', p.text.replace('—', ':'))

    apply_docx_layout(doc)
    doc.save(DOCX)


def set_font(run, name='Noto Sans CJK KR', size=None, bold=None, color=None):
    run.font.name = name
    run._element.rPr.rFonts.set(qn('w:eastAsia'), name)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.bold = bold
    if color is not None:
        run.font.color.rgb = RGBColor(*color)


def add_page_field(paragraph):
    run = paragraph.add_run()
    fld_char = OxmlElement('w:fldChar')
    fld_char.set(qn('w:fldCharType'), 'begin')
    instr = OxmlElement('w:instrText')
    instr.set(qn('xml:space'), 'preserve')
    instr.text = ' PAGE '
    sep = OxmlElement('w:fldChar')
    sep.set(qn('w:fldCharType'), 'separate')
    end = OxmlElement('w:fldChar')
    end.set(qn('w:fldCharType'), 'end')
    run._r.extend((fld_char, instr, sep, end))
    set_font(run, size=8, color=(90, 90, 90))


def apply_docx_layout(doc: Document) -> None:
    sec = doc.sections[0]
    sec.page_width = Mm(210)
    sec.page_height = Mm(297)
    sec.top_margin = Mm(18)
    sec.bottom_margin = Mm(19)
    sec.left_margin = Mm(20)
    sec.right_margin = Mm(20)
    sec.header_distance = Mm(8)
    sec.footer_distance = Mm(8)

    styles = doc.styles
    normal = styles['Normal']
    normal.font.name = 'Noto Sans CJK KR'
    normal._element.rPr.rFonts.set(qn('w:eastAsia'), 'Noto Sans CJK KR')
    normal.font.size = Pt(9.5)
    normal.paragraph_format.line_spacing = 1.35
    normal.paragraph_format.space_after = Pt(4)

    if 'Report Subtitle' not in styles:
        sub = styles.add_style('Report Subtitle', WD_STYLE_TYPE.PARAGRAPH)
    else:
        sub = styles['Report Subtitle']
    sub.font.name = 'Noto Sans CJK KR'
    sub._element.rPr.rFonts.set(qn('w:eastAsia'), 'Noto Sans CJK KR')
    sub.font.size = Pt(13)
    sub.font.bold = True
    sub.font.color.rgb = RGBColor(56, 78, 105)
    sub.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.paragraph_format.space_after = Pt(5)

    title = styles['Title']
    title.font.name = 'Noto Sans CJK KR'
    title._element.rPr.rFonts.set(qn('w:eastAsia'), 'Noto Sans CJK KR')
    title.font.size = Pt(24)
    title.font.bold = True
    title.font.color.rgb = RGBColor(25, 49, 76)
    title.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER

    for name, size, color in (('Heading 1', 15, (25, 49, 76)),
                              ('Heading 2', 12, (49, 78, 108)),
                              ('Heading 3', 10.5, (49, 78, 108))):
        st = styles[name]
        st.font.name = 'Noto Sans CJK KR'
        st._element.rPr.rFonts.set(qn('w:eastAsia'), 'Noto Sans CJK KR')
        st.font.size = Pt(size)
        st.font.bold = True
        st.font.color.rgb = RGBColor(*color)
        st.paragraph_format.space_before = Pt(10 if name == 'Heading 1' else 7)
        st.paragraph_format.space_after = Pt(4)
        st.paragraph_format.keep_with_next = True
        st.paragraph_format.keep_together = True

    doc.paragraphs[0].style = styles['Title']
    doc.paragraphs[1].style = sub
    doc.paragraphs[2].style = sub
    for i in (3, 4):
        doc.paragraphs[i].alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in doc.paragraphs[i].runs:
            set_font(run, size=8.5, color=(90, 90, 90))

    for i, p in enumerate(doc.paragraphs):
        text = p.text.strip()
        if i in (0, 1, 2, 3, 4):
            continue
        if text == '국문 초록' or re.match(r'^(?:\d+|부록)\.(?:\s|$)', text):
            p.style = styles['Heading 1']
        elif re.match(r'^\d+\.\d+\s', text):
            p.style = styles['Heading 2']
        elif re.match(r'^(그림|표)\s+\d+\.', text):
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.paragraph_format.keep_with_next = True
            p.paragraph_format.space_before = Pt(2)
            p.paragraph_format.space_after = Pt(4)
            for run in p.runs:
                set_font(run, size=8, color=(80, 80, 80))
        elif text.startswith('[') and re.match(r'^\[\d+\]', text):
            p.paragraph_format.left_indent = Mm(5)
            p.paragraph_format.first_line_indent = Mm(-5)
            p.paragraph_format.space_after = Pt(2)
            for run in p.runs:
                set_font(run, size=8)
        elif text:
            p.paragraph_format.widow_control = True
            p.paragraph_format.keep_together = False
            if not text.startswith(('주제어:', 'JEL 분류:')):
                p.paragraph_format.first_line_indent = Mm(4)
            for run in p.runs:
                set_font(run, size=9.5)

    # 제목 다음 문단과 빈 단락이 페이지 끝에 홀로 남지 않도록 한다.
    for p in doc.paragraphs:
        if p.style.name.startswith('Heading'):
            p.paragraph_format.keep_with_next = True

    # 표 머리행 반복, 행 분할 금지, 셀 여백과 글꼴 통일.
    for table in doc.tables:
        table.autofit = True
        header = table.rows[0]._tr
        tr_pr = header.get_or_add_trPr()
        tbl_header = OxmlElement('w:tblHeader')
        tbl_header.set(qn('w:val'), 'true')
        tr_pr.append(tbl_header)
        for row_idx, row in enumerate(table.rows):
            tr_pr = row._tr.get_or_add_trPr()
            cant_split = OxmlElement('w:cantSplit')
            tr_pr.append(cant_split)
            for cell in row.cells:
                tc_pr = cell._tc.get_or_add_tcPr()
                tc_mar = tc_pr.first_child_found_in('w:tcMar')
                if tc_mar is None:
                    tc_mar = OxmlElement('w:tcMar')
                    tc_pr.append(tc_mar)
                for side, val in (('top', 60), ('left', 70), ('bottom', 60), ('right', 70)):
                    node = tc_mar.find(qn(f'w:{side}'))
                    if node is None:
                        node = OxmlElement(f'w:{side}')
                        tc_mar.append(node)
                    node.set(qn('w:w'), str(val))
                    node.set(qn('w:type'), 'dxa')
                for p in cell.paragraphs:
                    p.paragraph_format.space_after = Pt(1)
                    p.paragraph_format.line_spacing = 1.05
                    for run in p.runs:
                        set_font(run, size=7.5, bold=(row_idx == 0))

    footer = sec.footer
    footer.is_linked_to_previous = False
    p = footer.paragraphs[0]
    p.clear()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_page_field(p)


def replace_para_text(p, old: str, new: str) -> bool:
    runs = list(p.runs)
    if not runs:
        return False
    joined = ''.join(r.text for r in runs)
    if old not in joined:
        return False
    joined = joined.replace(old, new)
    runs[0].text = joined
    for run in runs[1:]:
        run.text = ''
    return True


def replace_in_slide(slide, old: str, new: str) -> int:
    n = 0
    for shape in slide.shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        for p in shape.text_frame.paragraphs:
            if replace_para_text(p, old, new):
                n += 1
    return n


def make_failure_chart() -> dict[str, int]:
    aliases = {
        'Fitness': 'LOW_FITNESS',
        'Sharpe': 'LOW_SHARPE',
        'Sub-universe Sharpe': 'LOW_SUB_UNIVERSE_SHARPE',
        'Turnover': 'HIGH_TURNOVER',
        'Weight': 'CONCENTRATED_WEIGHT',
        'Weight concentration': 'CONCENTRATED_WEIGHT',
    }
    counts: Counter[str] = Counter()
    with sqlite3.connect(DB) as con:
        q = f"""SELECT fail_items FROM alphas
                WHERE ts <= ? AND sharpe IS NOT NULL
                  AND origin IN ({','.join('?' for _ in ORIGINS)})"""
        for (raw,) in con.execute(q, (CUTOFF, *ORIGINS)):
            try:
                items = json.loads(raw or '[]')
            except (TypeError, ValueError):
                continue
            for item in items:
                name = item.get('name', '') if isinstance(item, dict) else str(item)
                counts[aliases.get(name, name)] += 1

    labels = [x[0] for x in counts.most_common(10)][::-1]
    values = [counts[x] for x in labels]
    short = {
        'LOW_GLB_AMER_SHARPE': 'GLB_AMER',
        'LOW_GLB_APAC_SHARPE': 'GLB_APAC',
        'LOW_GLB_EMEA_SHARPE': 'GLB_EMEA',
        'LOW_SUB_UNIVERSE_SHARPE': 'SUB_UNIVERSE',
        'IS_LADDER_SHARPE': 'IS_LADDER',
        'LOW_2Y_SHARPE': 'LOW_2Y',
        'CONCENTRATED_WEIGHT': 'CONCENTRATED',
    }
    shown = [short.get(x, x) for x in labels]
    font_path = '/home/arcosium/.local/share/fonts/NotoSansKR-Regular.otf'
    font_prop = font_manager.FontProperties(fname=font_path)
    plt.rcParams.update({'font.family': font_prop.get_name(), 'font.size': 13})
    fig, ax = plt.subplots(figsize=(14.4, 6.8), dpi=100)
    bars = ax.barh(shown, values, color='#8EA9C5', height=0.62)
    ax.bar_label(bars, labels=[f'{v:,}' for v in values], padding=8, fontsize=12,
                 color='#666666', fontproperties=font_prop)
    ax.set_xlabel('FAIL 체크 건수 (중복 허용)', color='#666666', fontproperties=font_prop)
    for label in ax.get_yticklabels() + ax.get_xticklabels():
        label.set_fontproperties(font_prop)
    ax.grid(axis='x', color='#D8D8D8', linewidth=0.8)
    ax.set_axisbelow(True)
    ax.spines[['top', 'right']].set_visible(False)
    ax.tick_params(colors='#666666')
    ax.set_xlim(0, max(values) * 1.22)
    fig.tight_layout()
    fig.savefig(FAIL_FIG, transparent=False, facecolor='white')
    plt.close(fig)
    return dict(counts)


def ppt_edits() -> None:
    make_failure_chart()
    prs = Presentation(PPTX)

    edits = {
        3: [
            ('데이터필드 4,593 × 변환 6 × 결합 7 × 중립화 9 × 룩백 8 …',
             'USA/D1/TOP3000 Matrix 4,593 × 변환 6 × 결합 7 × 중립화 11 × 룩백 8 …'),
        ],
        4: [
            ('해당 3주(W30·W31·W32) 제출은 각각 15 · 12 · 10건이며,',
             '리더보드 집계에서 해당 3주(W30·W31·W32) 제출은 각각 15 · 12 · 10건이며,'),
            ('실제로 8월 13~14일에 재고 134건을 전수 발사했을 때 prod 상관 중앙값은 0.9028, 최소가 0.7111로 컷 0.7 아래가 한 건도 없었습니다.',
             '8월 13~14일 원장의 `[13시전 수동]`·`[13시전 블라스트]` 134건에서 prod 상관값은 53건 기록됐습니다. 중앙값은 0.9028, 최소는 0.7111로 컷 0.7 아래가 없었습니다.'),
        ],
        6: [
            ('4,593종 (Matrix)', '4,593종 (USA/D1/TOP3000 Matrix)'),
            ('±1 / TOP200~3000 / 0.08 고정', '±1 / TOP200~3000·TOPSP500 / 0.08 고정'),
            ('유전체 → FASTEXPR 렌더링 (실제 제출작 α#40383)', '유전체 → FASTEXPR 렌더링 (실제 제출작 α#40495, S 1.44)'),
            ('rank(-1*(vector_neut(ts_rank(winsorize(ts_backfill(bid_price_last_interval,120),std=4),3),\n                     ts_delta(winsorize(ts_backfill(high_price_30m_pre_close_2,120),std=4),10))))',
             'rank(group_rank(-1*(vector_neut(ts_rank(winsorize(ts_backfill(bid_price_last_interval, 120), std=4),5),\n                     ts_rank(winsorize(ts_backfill(high_price_30m_pre_close_2, 120), std=4),10))),subindustry))'),
        ],
        9: [
            ('2026-08-13 실측으로 복원한 게이트 프로파일', '2026-08-14 09:00 실측으로 복원한 게이트 프로파일'),
        ],
        12: [
            ('그럼에도 남아 있는 쏠림 (제출 성공작 25건의 패밀리)', '그럼에도 남아 있는 쏠림 (family 정보가 남은 제출작 25건)'),
            ('다목적 선택은 절반만 성공했습니다', '선택 효과와 생성 효과를 원장만으로 분리할 수 없습니다'),
            ('패밀리 엔트로피는 올랐지만 제출까지 간 알파는 여전히 model·pv 두 계열이 88%(22/25)입니다.',
             '제출작은 model 14건·pv 8건·imbalance 2건·news 1건이며, model·pv 두 계열이 88%(22/25)입니다.'),
            ('원인은 선택이 아니라 팔레트에 있었습니다.', '생성 팔레트의 쏠림도 함께 확인됐습니다.'),
            ('즉 다양성은 선택 단계가 아니라 생성 단계에서 만들어야 한다는 것이 이 실험의 결론입니다.',
             '따라서 다음 비교에서는 생성 팔레트와 엘리트 풀의 분포를 함께 기록해야 합니다.'),
        ],
        13: [
            ('35,701건 중 전 체크 통과 442건 = 1.2%', '35,701건 중 Sharpe 기록·fail_count 0·error_count 0인 442건 = 1.2%'),
        ],
        15: [
            ('S ≥ 1.58(표준 제출컷)도달 비율은', 'Sharpe가 기록된 주요 자동 생성 경로 6개(n=13,298)에서 S ≥ 1.58 도달 비율은'),
        ],
        16: [
            ('35,701건의 판정에서 뽑은 차단 사유입니다.', 'Sharpe가 기록된 주요 자동 생성 경로 6개(n=13,298)의 FAIL 체크를 중복 허용해 센 값입니다.'),
            ('① 상위 두 개(LOW_FITNESS · LOW_SHARPE)는 신호 세기 문제입니다. 데이터필드와 변환을 바꿔야 풀립니다.\n\n② 3~5위는 전부 리전별 Sharpe 입니다. GLB 알파는 세 리전을 동시에 만족해야 하므로 관문이 3개 더 붙습니다.\n\n③ 반면 PROD_CORRELATION은 124건으로 누적 순위가 낮습니다. 최근에는 이것이 실질 병목이 됐습니다. 누적 분포와 현재 병목은 다릅니다.',
             '① 상위 두 개(LOW_FITNESS · LOW_SHARPE)는 신호 세기 문제입니다.\n\n② LOW_2Y·IS_LADDER·SUB_UNIVERSE와 리전별 Sharpe도 반복됐습니다. GLB 알파는 확인할 관문이 많습니다.\n\n③ PROD_CORRELATION은 시뮬 단계 fail_items에는 18건뿐입니다. 제출 직전에 주로 관측되므로 누적 빈도가 작아도 최근 병목이 될 수 있습니다.'),
            ('8월 13~14일 재고 134건을 전수 발사했을 때 prod 상관 중앙값은 0.9028, 0.7 미만은 0건이었습니다.',
             '8월 13~14일 원장 태그 134건에서 prod 상관값은 53건 기록됐고 중앙값은 0.9028, 0.7 미만은 0건이었습니다.'),
        ],
        17: [
            ('Sharpe만 좇던 시기에 엘리트 풀이 rsk70 계열 한 가계로 수렴했습니다. 제출 21건 중 11건이 같은 데이터셋이었습니다.',
             '2026년 7월 28일~8월 3일 제출 21건 중 11건이 rsk70 데이터셋에서 나왔습니다.'),
            ('원인은 선택 알고리즘이 아니라 시딩 축 목록이 전부 rsk70이었다는 것이었습니다.',
             '당시 시딩 축 목록 26개가 모두 rsk70이어서 생성 단계부터 후보가 치우쳤습니다.'),
            ('NSGA-II와 최근성 윈도우로 풀은 다시 순환했지만, 생성 팔레트를 고치기 전까지 제출작 분포는 움직이지 않았습니다.',
             'NSGA-II와 최근성 윈도우를 함께 적용했지만 선택 효과만 따로 떼어 확인할 지표는 없었습니다.'),
        ],
        18: [
            ('제출 21건 중 11건이 rsk70 한 데이터셋이었고,', '2026년 7월 28일~8월 3일 제출 21건 중 11건이 rsk70 한 데이터셋이었고,'),
        ],
        19: [
            ('GLB/D1 기준 Broker·Imbalance는 필드 0이라 목표에서 제외했습니다.', '8월 14일 생성 대상 스냅샷에서 GLB/D1 Broker·Imbalance는 필드 0이라 목표에서 제외했습니다.'),
            ('미개척 카테고리(Analyst 4,218필드 · Fundamental 4,428필드)', '같은 스냅샷의 미개척 카테고리(Analyst 4,218필드 · Fundamental 4,428필드)'),
        ],
    }

    for slide_no, pairs in edits.items():
        slide = prs.slides[slide_no - 1]
        for old, new in pairs:
            replace_in_slide(slide, old, new)

    # 모든 슬라이드에서 산문용 긴 대시를 제거하고 붙은 문장을 띄운다.
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            for p in shape.text_frame.paragraphs:
                runs = list(p.runs)
                if not runs:
                    continue
                text = ''.join(r.text for r in runs)
                fixed = re.sub(r'\.(?=[가-힣])', '. ', text.replace('—', ':'))
                if fixed != text:
                    runs[0].text = fixed
                    for run in runs[1:]:
                        run.text = ''

    # fig5_failures 이미지를 갱신된 파일로 교체한다.
    slide = prs.slides[15]
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            # 이 슬라이드의 큰 그림 하나가 실패 원인 그래프다.
            if shape.width > 4_000_000:
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                sp = shape._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(str(FAIL_FIG), left, top, width, height)
                break

    prs.save(PPTX)


if __name__ == '__main__':
    report_edits()
    ppt_edits()
