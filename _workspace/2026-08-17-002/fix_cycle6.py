from pathlib import Path

from docx import Document
from pptx import Presentation


ROOT = Path('/home/arcosium/projects/GenomicWQB')
DOCX = ROOT / 'docs/유전알고리즘_알파리서치_리포트.docx'
PPTX = ROOT / 'docs/머신발표/GenomicWQB_머신발표.pptx'


def main():
    doc = Document(DOCX)
    for p in doc.paragraphs:
        for run in p.runs:
            run.text = run.text.replace('`', '')
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    for run in p.runs:
                        run.text = run.text.replace('`', '')
    doc.save(DOCX)

    prs = Presentation(PPTX)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.text = run.text.replace('`', '')
                    run.text = run.text.replace('같은 스냅샷의 같은 스냅샷의', '같은 스냅샷의')
    prs.save(PPTX)


if __name__ == '__main__':
    main()
