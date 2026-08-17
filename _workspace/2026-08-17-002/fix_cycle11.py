import re
from pathlib import Path

from docx import Document
from pptx import Presentation


ROOT = Path('/home/arcosium/projects/GenomicWQB')
DOCX = ROOT / 'docs/유전알고리즘_알파리서치_리포트.docx'
PPTX = ROOT / 'docs/머신발표/GenomicWQB_머신발표.pptx'
CONNECTIVE_COMMA = re.compile(r'(고|며|지만|아서|어서),')


def main():
    doc = Document(DOCX)
    doc.paragraphs[33].text = (
        '룩백을 여덟 종으로 제한한 것은 BRAIN의 권고를 그대로 따른 것이다[8]. '
        '37일이나 80일 같은 창은 데이터에 맞춘 흔적으로 읽히고 과적합 신호가 된다. '
        '표준 창만 허용하면 표현력은 줄지만 결과를 방어할 수 있다.'
    )
    for p in doc.paragraphs:
        for run in p.runs:
            fixed = CONNECTIVE_COMMA.sub(r'\1', run.text)
            if fixed != run.text:
                run.text = fixed
    doc.save(DOCX)

    prs = Presentation(PPTX)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    fixed = CONNECTIVE_COMMA.sub(r'\1', run.text)
                    if fixed != run.text:
                        run.text = fixed
    prs.save(PPTX)


if __name__ == '__main__':
    main()
