from pathlib import Path

from docx import Document
from pptx import Presentation


ROOT = Path('/home/arcosium/projects/GenomicWQB')
DOCX = ROOT / 'docs/유전알고리즘_알파리서치_리포트.docx'
PPTX = ROOT / 'docs/머신발표/GenomicWQB_머신발표.pptx'


def set_ppt_paragraph(paragraph, text):
    runs = list(paragraph.runs)
    if runs:
        runs[0].text = text
        for run in runs[1:]:
            run.text = ''
    else:
        paragraph.text = text


def fix_ppt():
    prs = Presentation(PPTX)

    # 실제 제출작 α#40495의 원장 code를 두 줄로 표시한다.
    formula_shape = prs.slides[5].shapes[68]
    set_ppt_paragraph(
        formula_shape.text_frame.paragraphs[0],
        'rank(group_rank(-1*(vector_neut(ts_rank(winsorize(ts_backfill(bid_price_last_interval, 120), std=4),5),',
    )
    set_ppt_paragraph(
        formula_shape.text_frame.paragraphs[1],
        '                     ts_rank(winsorize(ts_backfill(high_price_30m_pre_close_2, 120), std=4),10))),subindustry))',
    )

    # 새 FAIL 체크 그래프와 설명을 맞춘다.
    s16 = prs.slides[15]
    for shape in s16.shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        for p in shape.text_frame.paragraphs:
            if p.text.startswith('② 3~5위는'):
                set_ppt_paragraph(p, '② LOW_2Y·IS_LADDER·SUB_UNIVERSE와 리전별 Sharpe도 반복됐습니다. GLB 알파는 확인할 관문이 많습니다.')
            elif p.text.startswith('③ 반면 PROD_CORRELATION은'):
                set_ppt_paragraph(p, '③ PROD_CORRELATION은 시뮬 단계 fail_items에 18건뿐입니다. 제출 직전에 주로 관측되므로 누적 빈도가 작아도 최근 병목이 될 수 있습니다.')

    # 134건 태그 배치에는 성공작이 없었다. 별도 큐의 성공과 섞지 않는다.
    s18 = prs.slides[17]
    for shape in s18.shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        for p in shape.text_frame.paragraphs:
            if p.text.startswith('다음 날 재고를 전수 발사한 결과 134발'):
                set_ppt_paragraph(p, '다음 날 원장의 `[13시전 수동]`·`[13시전 블라스트]` 태그 134건 가운데 prod 상관값은 53건 기록됐습니다. 최소값은 0.7111로 컷 0.7 아래가 없었습니다.')
            elif p.text.startswith('2026년 7월 28일~8월 3일 2026년'):
                set_ppt_paragraph(p, p.text.replace('2026년 7월 28일~8월 3일 ', '', 1))

    prs.save(PPTX)


def remove_paragraph(paragraph):
    element = paragraph._element
    element.getparent().remove(element)
    paragraph._p = paragraph._element = None


def fix_docx():
    doc = Document(DOCX)

    # 인용되지 않고 본문 근거로도 쓰이지 않은 Fama–MacBeth 항목을 뺀다.
    for p in list(doc.paragraphs):
        if p.text.startswith('[7] Fama,'):
            remove_paragraph(p)
            break

    for p in doc.paragraphs:
        if '무작위 탐색은 2,582건 가운데' in p.text and '[6]' not in p.text:
            p.text = p.text + ' 무작위 탐색은 효율뿐 아니라 탐색 범위를 확보하기 위한 기준선이기도 하다[6].'
        p.text = p.text.replace('위험을 낮출 수 있다[9].', '위험을 낮출 수 있다[8].')
        p.text = p.text.replace('남는다[8].', '남는다[7].')
        if p.text.startswith('[8] Bailey,'):
            p.text = p.text.replace('[8]', '[7]', 1)
        elif p.text.startswith('[9] WorldQuant BRAIN.'):
            p.text = p.text.replace('[9]', '[8]', 1)
        elif p.text.startswith('[10] 2026 WorldQuant'):
            p.text = p.text.replace('[10]', '[9]', 1)
        elif p.text.startswith('[11] 분석 원장:'):
            p.text = p.text.replace('[11]', '[10]', 1)

    doc.save(DOCX)


if __name__ == '__main__':
    fix_ppt()
    fix_docx()
