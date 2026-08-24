#!/usr/bin/env python3
"""Build the 10-minute script document from the current presentation.

The presentation is the source of truth: visible headlines, presenter notes,
and presenter identity are read from the PPTX at runtime.
"""

from __future__ import annotations

import argparse
import hashlib
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.util import Inches

from build_single_axis_presentation import make_script_docx


TIMINGS = [30, 50, 50, 50, 70, 50, 75, 65, 55, 70, 35]
SECTION_TITLES = [
    "표지",
    "연구 질문",
    "검증 구조",
    "외부 통제 실험",
    "외부 실험 결과",
    "GenomicWQB 계보",
    "WQB 핵심 결과",
    "변이 폭별 결과",
    "강건성과 범위",
    "운용 규칙과 후속 검증",
    "결론",
]
TRANSITIONS = [
    "먼저 왜 변이 폭이 중요한지 보겠습니다.",
    "이 질문을 한 자료만으로 판단하지 않았습니다.",
    "외부 실험의 조건부터 말씀드리겠습니다.",
    "변이 폭만 바꿨을 때 결과는 분명한 경계로 나타났습니다.",
    "이제 이 원리가 퀀트 알파에서도 나타나는지 보겠습니다.",
    "세 방법의 결과는 거의 같은 크기로 모였습니다.",
    "다만 정확한 경계는 1축과 나머지 전체를 가르는 선이 아니었습니다.",
    "이 구분을 지키면 결론의 범위도 명확해집니다.",
    "이 결론을 실제 탐색 정책으로 옮기면 다음과 같습니다.",
    "마지막으로 세 문장만 남기겠습니다.",
    "발표를 마칩니다.",
]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--pptx", type=Path, required=True)
    parser.add_argument("--docx", type=Path, required=True)
    return parser.parse_args()


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def clean_text(text: str) -> str:
    return " ".join(text.split())


def visible_texts(slide):
    rows = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = clean_text(shape.text)
        if not text:
            continue
        rows.append((shape, text))
    return rows


def headline(slide, index: int) -> str:
    rows = visible_texts(slide)
    if index == 1:
        for _, text in rows:
            if "유전 알고리즘으로" in text:
                return text
    if index == 11:
        for _, text in rows:
            if "국소 탐색" in text:
                return text
    candidates = []
    for shape, text in rows:
        top = shape.top / Inches(1)
        width = shape.width / Inches(1)
        if 0.55 <= top <= 1.4 and width >= 5 and "/ 11" not in text:
            candidates.append((len(text), text))
    if not candidates:
        raise RuntimeError(f"slide {index}: headline not found")
    return max(candidates)[1]


def cover_message(slide) -> str:
    for _, text in visible_texts(slide):
        if "외부 통제 실험" in text and "실제 알파 계보" in text:
            return text
    return headline(slide, 1)


def presenter_identity(slide) -> str:
    for _, text in visible_texts(slide):
        if "소속" in text:
            return text
    return "발표자 정보는 표지와 동일"


def build_items(prs: Presentation) -> tuple[str, list[dict]]:
    if len(prs.slides) != 11:
        raise RuntimeError(f"expected 11 slides, found {len(prs.slides)}")
    identity = presenter_identity(prs.slides[0])
    items = []
    for index, slide in enumerate(prs.slides, 1):
        note = clean_text(slide.notes_slide.notes_text_frame.text)
        if not note:
            raise RuntimeError(f"slide {index}: presenter note is empty")
        message = cover_message(slide) if index == 1 else headline(slide, index)
        items.append(
            {
                "slide": index,
                "title": SECTION_TITLES[index - 1],
                "seconds": TIMINGS[index - 1],
                "message": message,
                "script": note,
                "transition": TRANSITIONS[index - 1],
            }
        )
    if sum(item["seconds"] for item in items) != 600:
        raise RuntimeError("presentation timing must total exactly 600 seconds")
    return identity, items


def update_meta(path: Path, digest: str) -> None:
    doc = Document(path)
    for paragraph in doc.paragraphs:
        if "원본 PDF SHA-256" in paragraph.text:
            paragraph.text = f"총 11장 · 10분 00초 · 기준 PPTX SHA-256 {digest}"
            paragraph.alignment = 1
        if "질문이 나오면 외부 실험은" in paragraph.text:
            element = paragraph._element
            element.getparent().remove(element)
    doc.save(path)


def main() -> None:
    args = parse_args()
    prs = Presentation(args.pptx)
    identity, items = build_items(prs)
    digest = sha256(args.pptx)
    make_script_docx(args.docx, identity, items, digest)
    update_meta(args.docx, digest)
    print(args.docx)


if __name__ == "__main__":
    main()
