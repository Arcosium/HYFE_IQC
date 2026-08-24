#!/usr/bin/env python3
"""Build the professional redesign of the single-axis mutation presentation.

The source report and generated deck live in the vault. Personal details are
read from the report PDF at runtime and are never embedded in this script.
"""

from __future__ import annotations

import argparse
import math
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from build_single_axis_presentation import load_report, slide_scripts


W, H = 13.333, 7.5
NAVY = RGBColor(8, 22, 34)
NAVY_2 = RGBColor(14, 39, 56)
PANEL = RGBColor(19, 49, 68)
OFFWHITE = RGBColor(244, 247, 248)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(18, 35, 46)
SLATE = RGBColor(91, 112, 127)
MUTED_DARK = RGBColor(153, 174, 187)
GRID = RGBColor(219, 229, 234)
TEAL = RGBColor(0, 169, 157)
TEAL_DARK = RGBColor(0, 118, 112)
TEAL_LIGHT = RGBColor(218, 242, 239)
AMBER = RGBColor(242, 162, 58)
AMBER_DARK = RGBColor(189, 110, 22)
AMBER_LIGHT = RGBColor(252, 235, 210)
CORAL = RGBColor(226, 100, 82)
BLUE = RGBColor(42, 122, 159)
BLUE_LIGHT = RGBColor(218, 235, 243)

FONT = "Pretendard"
MONO = "Noto Sans Mono CJK KR"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--source-pdf", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    return parser.parse_args()


def set_bg(slide, color) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, kind, x, y, w, h, *, fill, line=None, radius=True, width=1.0):
    if kind == "rect":
        shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    else:
        shape_kind = kind
    shape = slide.shapes.add_shape(shape_kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(width)
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=GRID, width=1.0, end_arrow=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if end_arrow:
        line.line.end_arrowhead = True
    return line


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=18,
    color=INK,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.0,
    spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_runs(slide, runs, x, y, w, h, *, size=18, color=INK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for spec in runs:
        run = p.add_run()
        run.text = spec[0]
        run.font.name = spec[4] if len(spec) > 4 else FONT
        run.font.size = Pt(spec[1] if len(spec) > 1 else size)
        run.font.bold = spec[2] if len(spec) > 2 else False
        run.font.color.rgb = spec[3] if len(spec) > 3 else color
    return box


def add_dark_grid(slide) -> None:
    for x in [0.55, 2.6, 4.65, 6.7, 8.75, 10.8, 12.78]:
        add_line(slide, x, 0.0, x, 7.5, color=RGBColor(17, 43, 59), width=0.55)
    for y in [0.7, 2.2, 3.7, 5.2, 6.7]:
        add_line(slide, 0.0, y, 13.333, y, color=RGBColor(17, 43, 59), width=0.55)


def header(slide, number, section, title, *, dark=False):
    ink = WHITE if dark else INK
    muted = MUTED_DARK if dark else SLATE
    add_text(slide, f"{number:02d}  /  11", 0.55, 0.32, 1.1, 0.25, size=9.5, color=TEAL, bold=True)
    add_text(slide, section.upper(), 1.75, 0.32, 3.0, 0.25, size=9.5, color=muted, bold=True)
    add_text(slide, title, 0.55, 0.76, 12.15, 0.55, size=25, color=ink, bold=True)
    add_shape(slide, "rect", 0.55, 1.42, 0.42, 0.04, fill=TEAL, radius=False)
    add_shape(slide, "rect", 0.97, 1.42, 11.8, 0.015, fill=RGBColor(52, 76, 91) if dark else GRID, radius=False)


def footer(slide, number, *, dark=False):
    color = MUTED_DARK if dark else SLATE
    add_text(slide, "GENOMICWQB  ·  EVOLUTIONARY ALPHA RESEARCH", 0.55, 7.14, 5.6, 0.18, size=7.5, color=color, bold=True)
    add_text(slide, f"{number:02d}", 12.35, 7.12, 0.42, 0.18, size=8, color=color, bold=True, align=PP_ALIGN.RIGHT)


def chip(slide, text, x, y, w, *, fill, color, line=None, size=10.5):
    add_shape(slide, "rect", x, y, w, 0.34, fill=fill, line=line, radius=True)
    add_text(slide, text, x, y + 0.015, w, 0.25, size=size, color=color, bold=True, align=PP_ALIGN.CENTER)


def metric(slide, label, value, x, y, w, *, accent=TEAL, dark=False, note=""):
    bg = PANEL if dark else WHITE
    line = RGBColor(40, 72, 90) if dark else GRID
    main = WHITE if dark else INK
    muted = MUTED_DARK if dark else SLATE
    add_shape(slide, "rect", x, y, w, 1.04, fill=bg, line=line, radius=True)
    add_text(slide, label, x + 0.15, y + 0.12, w - 0.3, 0.26, size=9.5, color=muted, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, value, x + 0.15, y + 0.39, w - 0.3, 0.42, size=23, color=main, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if note:
        add_text(slide, note, x + 0.15, y + 0.79, w - 0.3, 0.18, size=8.5, color=muted, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_note(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def dna(slide, x, y, w, h, *, dark=True, segments=11):
    c1, c2 = TEAL, AMBER
    faint = RGBColor(53, 86, 104) if dark else GRID
    points_a, points_b = [], []
    for i in range(segments):
        t = i / (segments - 1)
        px = x + t * w
        phase = math.sin(t * math.pi * 2.2)
        ay = y + h * (0.5 + 0.33 * phase)
        by = y + h * (0.5 - 0.33 * phase)
        points_a.append((px, ay))
        points_b.append((px, by))
        add_line(slide, px, ay, px, by, color=faint, width=1.1)
    for points, color in [(points_a, c1), (points_b, c2)]:
        for i in range(len(points) - 1):
            add_line(slide, *points[i], *points[i + 1], color=color, width=2.4)
        for px, py in points:
            add_shape(slide, MSO_SHAPE.OVAL, px - 0.055, py - 0.055, 0.11, 0.11, fill=color, radius=False)


def stage_label(slide, number, title, subtitle, x, y, w, *, accent, active=False):
    fill = accent if active else WHITE
    text_color = WHITE if active else INK
    line = accent if active else GRID
    add_shape(slide, "rect", x, y, w, 1.48, fill=fill, line=line, radius=True)
    add_text(slide, number, x + 0.18, y + 0.11, w - 0.36, 0.26, size=11, color=WHITE if active else accent, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, title, x + 0.18, y + 0.48, w - 0.36, 0.38, size=18, color=text_color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, subtitle, x + 0.18, y + 0.96, w - 0.36, 0.34, size=10.5, color=WHITE if active else SLATE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def interval_plot(slide, x, y, w, h, estimates, *, dark=True):
    axis_color = RGBColor(72, 101, 117) if dark else GRID
    label_color = MUTED_DARK if dark else SLATE
    main_color = WHITE if dark else INK
    max_x = 8.0
    for tick in [0, 2, 4, 6, 8]:
        tx = x + (tick / max_x) * w
        add_line(slide, tx, y, tx, y + h, color=axis_color, width=0.7)
        add_text(slide, str(tick), tx - 0.2, y + h + 0.06, 0.4, 0.2, size=8, color=label_color, align=PP_ALIGN.CENTER)
    for idx, (label, value, low, high) in enumerate(estimates):
        py = y + 0.35 + idx * (h - 0.7) / max(1, len(estimates) - 1)
        lx = x + low / max_x * w
        hx = x + high / max_x * w
        vx = x + value / max_x * w
        add_text(slide, label, x - 1.7, py - 0.13, 1.5, 0.26, size=10, color=main_color, bold=True, align=PP_ALIGN.RIGHT)
        add_line(slide, lx, py, hx, py, color=MUTED_DARK if dark else BLUE, width=1.8)
        add_line(slide, lx, py - 0.09, lx, py + 0.09, color=MUTED_DARK if dark else BLUE, width=1.1)
        add_line(slide, hx, py - 0.09, hx, py + 0.09, color=MUTED_DARK if dark else BLUE, width=1.1)
        add_shape(slide, MSO_SHAPE.OVAL, vx - 0.075, py - 0.075, 0.15, 0.15, fill=TEAL, radius=False)
        add_text(slide, f"+{value:.2f}", hx + 0.12, py - 0.13, 0.8, 0.26, size=10, color=TEAL, bold=True)


def make_deck(identity: str, scripts: list[dict]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # 01 — Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_dark_grid(slide)
    chip(slide, "10 MIN RESEARCH BRIEF", 0.65, 0.55, 2.2, fill=PANEL, color=TEAL)
    add_text(slide, "유전 알고리즘으로\n알파 리서치하기", 0.65, 1.28, 7.2, 1.55, size=36, color=WHITE, bold=True, spacing=0.9)
    add_text(slide, "1축 변이는 다축 변이보다 효과적인가", 0.68, 3.08, 6.7, 0.45, size=20, color=MUTED_DARK, bold=True)
    add_runs(slide, [("PBO ", 12, True, MUTED_DARK), ("3,000쌍", 17, True, TEAL), ("   ×   WQB ", 12, True, MUTED_DARK), ("9,839쌍", 17, True, AMBER)], 0.68, 4.02, 6.6, 0.38)
    add_text(slide, "외부 통제 실험에서 원리를 확인하고 실제 알파 계보에서 다시 검증했다", 0.68, 4.58, 6.8, 0.5, size=14, color=WHITE)
    add_text(slide, identity, 0.68, 6.45, 7.4, 0.3, size=10.5, color=MUTED_DARK)
    dna(slide, 7.65, 1.1, 5.1, 4.9, dark=True, segments=13)
    add_text(slide, "LOCAL EXPLOITATION", 8.8, 6.22, 3.35, 0.22, size=8.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "CONDITIONAL ESCAPE", 8.8, 6.52, 3.35, 0.22, size=8.5, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_note(slide, scripts[0]["script"])

    # 02 — Research question
    slide = prs.slides.add_slide(blank)
    set_bg(slide, OFFWHITE)
    header(slide, 2, "Research question", "변이 폭은 탐색 예산의 방향을 결정한다")
    chip(slide, "EXPLOIT", 0.7, 1.75, 1.15, fill=TEAL_LIGHT, color=TEAL_DARK)
    chip(slide, "EXPLORE", 11.45, 1.75, 1.15, fill=AMBER_LIGHT, color=AMBER_DARK)
    add_text(slide, "가까이, 정확하게", 0.72, 2.28, 3.25, 0.42, size=24, color=TEAL_DARK, bold=True)
    add_text(slide, "멀리, 한 번에", 9.35, 2.28, 3.25, 0.42, size=24, color=AMBER_DARK, bold=True, align=PP_ALIGN.RIGHT)
    genes = ["FIELD", "OP", "LOOKBACK", "NEUT", "DECAY", "UNIVERSE"]
    gx = 2.2
    for i, gene in enumerate(genes):
        x = gx + i * 1.48
        fill = TEAL if i == 2 else WHITE
        color = WHITE if i == 2 else SLATE
        add_shape(slide, "rect", x, 3.22, 1.12, 0.58, fill=fill, line=TEAL if i == 2 else GRID, radius=True)
        add_text(slide, gene, x, 3.35, 1.12, 0.28, size=11.25, color=color, bold=True, font=MONO, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(genes) - 1:
            add_line(slide, x + 1.12, 3.51, x + 1.48, 3.51, color=GRID, width=1.2)
    add_text(slide, "1축", 1.05, 3.28, 0.85, 0.4, size=18, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "다축", 11.45, 3.28, 0.85, 0.4, size=18, color=AMBER_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, 1.55, 4.85, 11.75, 4.85, color=GRID, width=1.5)
    add_text(slide, "구조 보존  ·  원인 추적  ·  좋은 부모 주변", 0.72, 5.08, 5.1, 0.35, size=12, color=SLATE)
    add_text(slide, "새 영역 탐색  ·  정체 탈출  ·  축 간 상호작용", 7.35, 5.08, 5.25, 0.35, size=12, color=SLATE, align=PP_ALIGN.RIGHT)
    add_shape(slide, "rect", 1.42, 5.82, 10.5, 0.84, fill=NAVY, radius=True)
    add_text(slide, "좋은 후보를 다듬을 때, 1축 변이가 더 효과적인가?", 1.7, 6.05, 9.95, 0.38, size=21, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 2)
    add_note(slide, scripts[1]["script"])

    # 03 — Evidence architecture
    slide = prs.slides.add_slide(blank)
    set_bg(slide, OFFWHITE)
    header(slide, 3, "Evidence architecture", "서로 다른 자료에서 같은 방향이 나와야 결론이 단단해진다")
    stages = [
        ("01", "이론", "좁은 변이가 국소 탐색에 유리", TEAL),
        ("02", "통제 실험", "PBO 3,000쌍 · 변이 폭만 통제", TEAL),
        ("03", "WQB 실증", "실평가 계보 9,839쌍", AMBER),
        ("04", "운용 원칙", "활용 1~2축 · 정체 시 3축+", NAVY_2),
    ]
    xs = [0.65, 3.85, 7.05, 10.25]
    for i, ((num, title, sub, accent), x) in enumerate(zip(stages, xs)):
        stage_label(slide, num, title, sub, x, 2.0, 2.45, accent=accent, active=(i == 3))
        if i < 3:
            add_line(slide, x + 2.47, 2.74, xs[i + 1] - 0.08, 2.74, color=SLATE, width=1.4, end_arrow=True)
    add_text(slide, "원리", 1.0, 4.15, 1.05, 0.25, size=10, color=TEAL_DARK, bold=True)
    add_line(slide, 1.95, 4.28, 6.35, 4.28, color=TEAL, width=2.2)
    add_shape(slide, MSO_SHAPE.OVAL, 6.25, 4.18, 0.2, 0.2, fill=TEAL, radius=False)
    add_text(slide, "현실 적용", 11.15, 4.15, 1.15, 0.25, size=10, color=AMBER_DARK, bold=True, align=PP_ALIGN.RIGHT)
    add_line(slide, 6.35, 4.28, 11.2, 4.28, color=AMBER, width=2.2)
    add_shape(slide, MSO_SHAPE.OVAL, 6.25, 4.18, 0.2, 0.2, fill=AMBER, radius=False)
    add_shape(slide, "rect", 2.0, 5.15, 9.3, 0.92, fill=WHITE, line=GRID, radius=True)
    add_text(slide, "이론 → 공개 문제 → 실제 알파 계보 → 상태별 정책", 2.25, 5.43, 8.8, 0.34, size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    chip(slide, "MECHANISM", 2.35, 6.28, 1.35, fill=TEAL_LIGHT, color=TEAL_DARK)
    chip(slide, "EVIDENCE", 5.05, 6.28, 1.25, fill=BLUE_LIGHT, color=BLUE)
    chip(slide, "VALIDATION", 7.65, 6.28, 1.4, fill=AMBER_LIGHT, color=AMBER_DARK)
    chip(slide, "POLICY", 10.35, 6.28, 1.05, fill=RGBColor(225, 232, 236), color=NAVY_2)
    footer(slide, 3)
    add_note(slide, scripts[2]["script"])

    # 04 — Controlled experiment
    slide = prs.slides.add_slide(blank)
    set_bg(slide, OFFWHITE)
    header(slide, 4, "Controlled experiment", "같은 시작점과 예산에서 바꾸는 축 수만 달리했다")
    metrics = [("PBO 문제군", "25종"), ("대응 비교", "3,000쌍"), ("알고리즘 실행", "6,000회"), ("함수평가", "696만 회")]
    for i, (label, value) in enumerate(metrics):
        metric(slide, label, value, 0.62 + i * 3.1, 1.72, 2.72, accent=TEAL if i < 2 else AMBER)
    add_shape(slide, "rect", 0.9, 3.18, 5.45, 2.14, fill=TEAL_LIGHT, line=TEAL, radius=True)
    chip(slide, "SINGLE-AXIS", 2.9, 3.48, 1.45, fill=TEAL, color=WHITE)
    add_text(slide, "매 평가마다 한 자리만 변경", 1.2, 4.03, 4.85, 0.42, size=22, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "부모 구조를 최대한 보존", 1.2, 4.6, 4.85, 0.3, size=12, color=SLATE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_shape(slide, "rect", 6.98, 3.18, 5.45, 2.14, fill=AMBER_LIGHT, line=AMBER, radius=True)
    chip(slide, "MULTI-AXIS", 8.98, 3.48, 1.45, fill=AMBER, color=NAVY)
    add_text(slide, "WQB의 2~14축 비율대로 변경", 7.28, 4.03, 4.85, 0.42, size=22, color=AMBER_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "실제 탐색의 변이 폭을 재현", 7.28, 4.6, 4.85, 0.3, size=12, color=SLATE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "동일", 6.28, 3.72, 0.76, 0.22, size=8.5, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "=", 6.3, 4.06, 0.72, 0.46, size=28, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "20d", 6.28, 4.7, 0.76, 0.22, size=8.5, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, "rect", 1.5, 5.82, 10.3, 0.82, fill=NAVY, radius=True)
    add_runs(slide, [("핵심  ", 12, True, TEAL), ("최적점 도달률", 17, True, WHITE), ("     보조  ", 12, True, AMBER), ("최적점까지 줄인 거리", 17, True, WHITE)], 1.8, 6.06, 9.7, 0.36, align=PP_ALIGN.CENTER)
    footer(slide, 4)
    add_note(slide, scripts[3]["script"])

    # 05 — External result
    slide = prs.slides.add_slide(blank)
    set_bg(slide, OFFWHITE)
    header(slide, 5, "External result", "국소 수렴은 1축이 강했고, 장벽 탈출은 다축이 앞섰다")
    add_text(slide, "최적점 도달률", 0.65, 1.72, 2.2, 0.3, size=12, color=SLATE, bold=True)
    add_text(slide, "45.5%", 0.65, 2.08, 2.35, 0.68, size=36, color=TEAL, bold=True)
    add_text(slide, "1축", 3.2, 2.25, 0.65, 0.3, size=12, color=INK, bold=True)
    add_shape(slide, "rect", 3.92, 2.19, 4.15, 0.38, fill=TEAL, radius=True)
    add_text(slide, "26.5%", 0.65, 3.08, 2.35, 0.68, size=36, color=AMBER, bold=True)
    add_text(slide, "다축", 3.2, 3.25, 0.65, 0.3, size=12, color=INK, bold=True)
    add_shape(slide, "rect", 3.92, 3.19, 2.42, 0.38, fill=AMBER, radius=True)
    add_shape(slide, "rect", 1.02, 4.23, 6.42, 1.2, fill=WHITE, line=GRID, radius=True)
    add_text(slide, "+19.0%p", 1.25, 4.36, 5.96, 0.48, size=31, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "95% CI  [8.8, 29.1]", 1.25, 4.86, 5.96, 0.25, size=11.5, color=SLATE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_shape(slide, "rect", 8.5, 1.83, 4.18, 1.6, fill=TEAL_LIGHT, line=TEAL, radius=True)
    add_text(slide, "기본 문제", 8.8, 2.05, 3.58, 0.28, size=10, color=SLATE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "1축 도달률  +69.2%p", 8.8, 2.42, 3.58, 0.48, size=21, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_shape(slide, "rect", 8.5, 3.7, 4.18, 1.6, fill=AMBER_LIGHT, line=AMBER, radius=True)
    add_text(slide, "험준 · 상호작용", 8.8, 3.92, 3.58, 0.28, size=10, color=SLATE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "다축 이동 거리  +0.230", 8.8, 4.29, 3.58, 0.48, size=21, color=AMBER_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_line(slide, 0.68, 6.15, 12.65, 6.15, color=GRID, width=1.1)
    add_text(slide, "1축은 좋은 후보를 다듬는다. 다축은 장벽을 넘는다.", 0.68, 6.4, 11.95, 0.4, size=21, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 5)
    add_note(slide, scripts[4]["script"])

    # 06 — Data lineage
    slide = prs.slides.add_slide(blank)
    set_bg(slide, OFFWHITE)
    header(slide, 6, "GenomicWQB lineage", "실제로 새 평가를 돌린 부모·자식 9,839쌍만 남겼다")
    funnel = [("전체 알파", "39,506건", 10.9, BLUE), ("Sharpe 연결", "12,168쌍", 7.6, TEAL), ("실평가 계보", "9,839쌍", 5.55, NAVY_2)]
    for i, (label, value, width, accent) in enumerate(funnel):
        x = 0.72 + (10.9 - width) / 2
        y = 1.72 + i * 1.08
        add_shape(slide, "rect", x, y, width, 0.78, fill=WHITE if i < 2 else NAVY_2, line=accent, radius=True)
        add_text(slide, label, x + 0.25, y + 0.08, width - 0.5, 0.25, size=10.5, color=SLATE if i < 2 else MUTED_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, value, x + 0.25, y + 0.34, width - 0.5, 0.34, size=22, color=accent if i < 2 else WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < 2:
            add_line(slide, 6.17, y + 0.78, 6.17, y + 1.08, color=SLATE, width=1.4, end_arrow=True)
    add_line(slide, 6.17, 4.64, 6.17, 5.02, color=SLATE, width=1.4)
    add_line(slide, 3.6, 5.02, 8.75, 5.02, color=SLATE, width=1.4)
    add_line(slide, 3.6, 5.02, 3.6, 5.28, color=TEAL, width=1.8, end_arrow=True)
    add_line(slide, 8.75, 5.02, 8.75, 5.28, color=AMBER, width=1.8, end_arrow=True)
    metric(slide, "1축  ·  genes_changed = 1", "3,402쌍", 1.75, 5.35, 3.75, accent=TEAL)
    metric(slide, "다축  ·  genes_changed ≥ 2", "6,437쌍", 6.88, 5.35, 3.75, accent=AMBER)
    chip(slide, "단순 비교", 3.45, 6.63, 1.25, fill=TEAL_LIGHT, color=TEAL_DARK)
    chip(slide, "조건 보정", 5.42, 6.63, 1.25, fill=BLUE_LIGHT, color=BLUE)
    chip(slide, "같은 부모", 7.4, 6.63, 1.25, fill=AMBER_LIGHT, color=AMBER_DARK)
    footer(slide, 6)
    add_note(slide, scripts[5]["script"])

    # 07 — WQB main result
    slide = prs.slides.add_slide(blank)
    set_bg(slide, OFFWHITE)
    header(slide, 7, "WQB result", "1축의 부모 개선률은 다축보다 4~5%p 높았다")
    add_text(slide, "부모 대비 Sharpe 개선률", 0.65, 1.72, 3.0, 0.3, size=12, color=SLATE, bold=True)
    add_text(slide, "16.05%", 0.65, 2.15, 2.25, 0.58, size=34, color=TEAL, bold=True)
    add_text(slide, "1축", 2.95, 2.29, 0.65, 0.28, size=11, color=INK, bold=True)
    add_shape(slide, "rect", 3.55, 2.25, 3.2, 0.34, fill=TEAL, radius=True)
    add_text(slide, "11.96%", 0.65, 3.02, 2.25, 0.58, size=34, color=AMBER, bold=True)
    add_text(slide, "다축", 2.95, 3.16, 0.65, 0.28, size=11, color=INK, bold=True)
    add_shape(slide, "rect", 3.55, 3.12, 2.38, 0.34, fill=AMBER, radius=True)
    add_text(slide, "+4.09%p", 0.65, 4.15, 3.0, 0.65, size=38, color=INK, bold=True)
    add_text(slide, "95% CI  [1.64, 6.49]", 0.67, 4.9, 2.8, 0.3, size=11.5, color=SLATE, bold=True)
    add_text(slide, "효과 차이  (%p)", 8.0, 1.72, 2.1, 0.25, size=10, color=SLATE, bold=True)
    interval_plot(slide, 8.0, 2.08, 3.7, 2.75, [("단순 비교", 4.09, 1.64, 6.49), ("조건 보정", 4.77, 2.53, 7.01), ("같은 부모", 4.51, 2.14, 6.89)], dark=False)
    add_shape(slide, "rect", 7.55, 5.4, 4.72, 0.88, fill=TEAL_LIGHT, line=TEAL, radius=True)
    add_text(slide, "세 계산이 모두 같은 방향", 7.8, 5.56, 4.2, 0.48, size=18, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "관측 자료에서도 국소 탐색 우위가 반복됐다", 1.15, 6.35, 11.0, 0.36, size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 7)
    add_note(slide, scripts[6]["script"])

    # 08 — Dose response
    slide = prs.slides.add_slide(blank)
    set_bg(slide, OFFWHITE)
    header(slide, 8, "Mutation breadth", "성과의 경계는 1축 대 다축이 아니라 1~2축 대 3축 이상이다")
    labels = ["1", "2", "3", "4", "5", "6+"]
    values = [16.05, 16.73, 13.13, 10.15, 6.03, 6.22]
    counts = [3402, 2158, 1462, 1103, 813, 901]
    medians = [-0.58, -0.51, -0.56, -0.81, -0.92, -0.97]
    x0, y0, cw, ch = 0.78, 1.92, 8.3, 3.8
    for tick in [0, 5, 10, 15, 20]:
        py = y0 + ch - tick / 20 * ch
        add_line(slide, x0, py, x0 + cw, py, color=GRID, width=0.7)
        add_text(slide, f"{tick}%", x0 - 0.55, py - 0.12, 0.45, 0.22, size=8.5, color=SLATE, align=PP_ALIGN.RIGHT)
    pts = []
    for i, (label, value, count) in enumerate(zip(labels, values, counts)):
        px = x0 + 0.55 + i * (cw - 1.1) / 5
        py = y0 + ch - value / 20 * ch
        pts.append((px, py))
        if i < 2:
            add_shape(slide, "rect", px - 0.42, y0, 0.84, ch, fill=TEAL_LIGHT, radius=False)
        add_text(slide, f"n={count:,}", px - 0.55, py - 0.43, 1.1, 0.22, size=8, color=SLATE, align=PP_ALIGN.CENTER)
        add_text(slide, label, px - 0.2, y0 + ch + 0.14, 0.4, 0.22, size=9.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
    for i in range(len(pts) - 1):
        add_line(slide, *pts[i], *pts[i + 1], color=BLUE, width=2.5)
    for px, py in pts:
        add_shape(slide, MSO_SHAPE.OVAL, px - 0.075, py - 0.075, 0.15, 0.15, fill=BLUE, radius=False)
    add_text(slide, "부모 대비 개선률", x0, y0 - 0.38, 2.3, 0.25, size=10.5, color=SLATE, bold=True)
    add_text(slide, "동시에 바꾼 유전자 수", x0 + 2.7, y0 + ch + 0.52, 3.0, 0.25, size=10.5, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_shape(slide, "rect", 9.55, 1.9, 3.0, 1.18, fill=TEAL_LIGHT, line=TEAL, radius=True)
    add_text(slide, "1~2축", 9.83, 2.08, 2.44, 0.28, size=12, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "16.05%  ·  16.73%", 9.83, 2.45, 2.44, 0.38, size=18, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_shape(slide, "rect", 9.55, 3.36, 3.0, 1.18, fill=AMBER_LIGHT, line=AMBER, radius=True)
    add_text(slide, "5축 이상", 9.83, 3.54, 2.44, 0.28, size=12, color=AMBER_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "약 6%", 9.83, 3.91, 2.44, 0.38, size=18, color=AMBER_DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_shape(slide, "rect", 9.55, 4.82, 3.0, 1.18, fill=WHITE, line=GRID, radius=True)
    add_text(slide, "ΔSharpe 중앙값", 9.83, 5.0, 2.44, 0.28, size=11, color=SLATE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "전 구간 음수", 9.83, 5.37, 2.44, 0.38, size=18, color=CORAL, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "절대 성공이 아니라 실패가 덜한 상대 우위다", 1.15, 6.42, 11.0, 0.35, size=18, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 8)
    add_note(slide, scripts[7]["script"])

    # 09 — Scope and robustness
    slide = prs.slides.add_slide(blank)
    set_bg(slide, OFFWHITE)
    header(slide, 9, "Robustness & scope", "국소 개선 우위는 강건하지만 최종 성공까지 보장하지 않는다")
    estimates = [("단순 비교", "+4.09%p", "부모 군집 부트스트랩"), ("조건 보정", "+4.77%p", "시점 · 경로 · 시장 설정 통제"), ("같은 부모", "+4.51%p", "부모 품질 차이 축소")]
    for i, (label, value, note) in enumerate(estimates):
        metric(slide, label, value, 0.68 + i * 4.18, 1.76, 3.72, accent=TEAL, note=note)
    add_shape(slide, "rect", 0.68, 3.36, 5.9, 2.24, fill=NAVY_2, radius=True)
    chip(slide, "SUPPORTED", 3.02, 3.67, 1.22, fill=TEAL, color=WHITE)
    add_text(slide, "국소 탐색에서\n1~2축은 유효하다", 0.98, 4.15, 5.3, 0.92, size=26, color=WHITE, bold=True, spacing=0.9, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_shape(slide, "rect", 6.78, 3.36, 5.86, 2.24, fill=AMBER_LIGHT, line=AMBER, radius=True)
    chip(slide, "NOT SUPPORTED", 8.94, 3.67, 1.55, fill=AMBER, color=NAVY)
    add_text(slide, "1축이 항상 최고다\n최종 성공을 보장한다", 7.08, 4.15, 5.26, 0.92, size=24, color=AMBER_DARK, bold=True, spacing=0.9, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_shape(slide, "rect", 2.25, 6.0, 8.83, 0.58, fill=WHITE, line=GRID, radius=True)
    add_runs(slide, [("S≥1.58 신규 도달   ", 11, True, SLATE), ("1축 1.56%", 14, True, TEAL_DARK), ("   ·   ", 11, True, SLATE), ("다축 1.48%", 14, True, AMBER_DARK), ("   →   뚜렷한 차이 없음", 11, True, SLATE)], 2.48, 6.17, 8.36, 0.25, align=PP_ALIGN.CENTER)
    footer(slide, 9)
    add_note(slide, scripts[8]["script"])

    # 10 — Operating rule
    slide = prs.slides.add_slide(blank)
    set_bg(slide, OFFWHITE)
    header(slide, 10, "Operating rule", "활용과 탈출을 상태에 따라 나눈다")
    blocks = [
        (0.7, 2.05, 2.55, "유망한 부모", "최근 개선 지속", WHITE, BLUE, INK),
        (3.7, 2.05, 2.75, "1~2축 우선", "국소 활용", TEAL, TEAL, WHITE),
        (6.9, 2.05, 2.15, "정체?", "연속 실패", WHITE, SLATE, INK),
        (9.5, 2.05, 2.85, "3축 이상", "장벽 탈출", AMBER, AMBER, NAVY),
    ]
    for i, (x, y, w, title, note, fill, line, color) in enumerate(blocks):
        add_shape(slide, "rect", x, y, w, 1.42, fill=fill, line=line, radius=True)
        add_text(slide, title, x + 0.2, y + 0.35, w - 0.4, 0.4, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
        note_color = SLATE if fill == WHITE else NAVY_2 if fill == AMBER else TEAL_LIGHT
        add_text(slide, note, x + 0.2, y + 0.9, w - 0.4, 0.24, size=10.5, color=note_color, bold=True, align=PP_ALIGN.CENTER)
        if i < len(blocks) - 1:
            nx = blocks[i + 1][0]
            add_line(slide, x + w, y + 0.72, nx - 0.08, y + 0.72, color=TEAL if i == 0 else AMBER if i == 2 else SLATE, width=1.7, end_arrow=True)
    add_text(slide, "NO", 7.77, 3.62, 0.45, 0.2, size=8.5, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, 7.98, 3.47, 7.98, 4.06, color=TEAL, width=1.4)
    add_line(slide, 7.98, 4.06, 5.08, 4.06, color=TEAL, width=1.4)
    add_line(slide, 5.08, 4.06, 5.08, 3.5, color=TEAL, width=1.4, end_arrow=True)
    add_shape(slide, "rect", 0.96, 4.72, 11.4, 1.28, fill=NAVY_2, radius=True)
    add_text(slide, "다음 검증", 1.28, 4.86, 1.35, 0.34, size=11, color=TEAL, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "같은 부모를 k=1 · k=2 · k≥3으로 무작위 배정", 2.58, 4.86, 6.6, 0.52, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "방법마다 최소\n약 1,176쌍", 9.5, 4.83, 2.25, 0.66, size=15, color=AMBER, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "관측 결론", 3.98, 6.4, 1.55, 0.28, size=12, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, 5.65, 6.54, 7.47, 6.54, color=TEAL, width=2.0, end_arrow=True)
    add_text(slide, "인과 결론", 7.62, 6.4, 1.55, 0.28, size=12, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, 10)
    add_note(slide, scripts[9]["script"])

    # 11 — Conclusion
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    add_dark_grid(slide)
    chip(slide, "CONCLUSION", 0.65, 0.58, 1.35, fill=PANEL, color=TEAL)
    add_text(slide, "1축 변이는 퀀트 알파의\n국소 탐색에서 유효하다", 0.65, 1.2, 7.8, 1.3, size=33, color=WHITE, bold=True, spacing=0.9)
    takeaways = [
        ("01", "외부 통제 실험과 WQB 계보가 같은 방향을 보였다", TEAL),
        ("02", "정확한 운용 경계는 1~2축 대 3축 이상이다", TEAL),
        ("03", "다축은 정체를 벗어나는 조건부 연산자다", AMBER),
    ]
    for i, (num, text, color) in enumerate(takeaways):
        y = 3.2 + i * 0.82
        add_shape(slide, MSO_SHAPE.OVAL, 0.7, y, 0.34, 0.34, fill=color, radius=False)
        add_text(slide, num, 0.7, y + 0.07, 0.34, 0.14, size=7, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, text, 1.28, y - 0.02, 7.25, 0.38, size=17, color=WHITE, bold=True)
    dna(slide, 8.42, 1.05, 4.2, 4.65, dark=True, segments=11)
    add_shape(slide, "rect", 9.22, 5.92, 2.62, 0.82, fill=PANEL, line=RGBColor(44, 78, 96), radius=True)
    add_text(slide, "Q&A", 9.22, 6.12, 2.62, 0.37, size=24, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "EXPLOIT LOCALLY  ·  ESCAPE CONDITIONALLY", 0.67, 6.68, 7.0, 0.22, size=8.5, color=MUTED_DARK, bold=True)
    add_note(slide, scripts[10]["script"])

    return prs


def main() -> None:
    args = parse_args()
    identity, _, _ = load_report(args.source_pdf)
    identity = re.sub(r"\s*·\s*교신\s*:\s*[^·]+$", "", identity).strip()
    scripts = slide_scripts()
    deck = make_deck(identity, scripts)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    deck.save(args.output)
    print(args.output)


if __name__ == "__main__":
    main()
