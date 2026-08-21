#!/usr/bin/env python3
"""Build a 10-minute presentation from the current alpha-research report PDF.

Personal details are read from the source PDF at runtime. Outputs are expected
to point into the vault; no identity value is embedded in this repository.
"""

from __future__ import annotations

import argparse
import hashlib
import re
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Mm, Pt, RGBColor as DocxRGBColor
from PIL import Image
from pypdf import PdfReader
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt as PptPt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(248, 247, 243)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(28, 34, 42)
MUTED = RGBColor(92, 101, 111)
LINE = RGBColor(214, 217, 220)
BLUE = RGBColor(42, 95, 137)
BLUE_LIGHT = RGBColor(222, 234, 242)
ORANGE = RGBColor(207, 121, 29)
ORANGE_LIGHT = RGBColor(247, 230, 207)
RED = RGBColor(160, 55, 47)
GREEN = RGBColor(54, 117, 84)
GRAY_CARD = RGBColor(238, 238, 235)

TITLE_FONT = "Noto Serif CJK KR"
BODY_FONT = "Noto Sans CJK KR"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--source-pdf", type=Path, required=True)
    parser.add_argument("--asset-dir", type=Path, required=True)
    parser.add_argument("--pptx", type=Path, required=True)
    parser.add_argument("--script-docx", type=Path, required=True)
    parser.add_argument("--audit-text", type=Path, required=True)
    return parser.parse_args()


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def load_report(path: Path) -> tuple[str, str, str]:
    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    full_text = "\n".join(pages)
    first = pages[0] if pages else ""
    identity_match = re.search(r"저자:\s*([^\n]+)", first)
    identity = identity_match.group(1).strip() if identity_match else "저자 정보는 원문과 동일"
    required = [
        "3,000", "45.5%", "26.5%", "39,506", "9,839",
        "16.05%", "11.96%", "4.09", "4.77", "4.51",
    ]
    missing = [token for token in required if token not in full_text]
    if missing:
        raise RuntimeError(f"source PDF is missing required values: {missing}")
    return identity, full_text, sha256(path)


def set_slide_bg(slide, color=BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_run(run, *, size=20, bold=False, color=INK, font=BODY_FONT, italic=False) -> None:
    run.font.name = font
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size=20,
    bold=False,
    color=INK,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.02,
    italic=False,
    line_spacing=1.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font, italic=italic)
    return box


def add_rich_lines(slide, lines, x, y, w, h, *, bullet=False, size=20, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for index, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.space_after = PptPt(9)
        paragraph.line_spacing = 1.05
        if bullet:
            paragraph.text = "•  " + line
        for run in paragraph.runs:
            set_run(run, size=size, color=color)
    return box


def add_rect(slide, x, y, w, h, *, fill=WHITE, line=LINE, radius=True, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    shape.line.color.rgb = line
    shape.line.width = PptPt(1)
    return shape


def add_card(slide, x, y, w, h, title, value, note="", *, accent=BLUE, value_size=28):
    add_rect(slide, x, y, w, h, fill=WHITE, line=LINE)
    add_rect(slide, x, y, 0.08, h, fill=accent, line=accent, radius=False)
    add_text(slide, title, x + 0.22, y + 0.17, w - 0.35, 0.35, size=12, bold=True, color=MUTED)
    add_text(slide, value, x + 0.22, y + 0.55, w - 0.35, 0.55, size=value_size, bold=True, color=accent)
    if note:
        add_text(slide, note, x + 0.22, y + h - 0.42, w - 0.35, 0.26, size=10.5, color=MUTED)


def add_title(slide, number: int, title: str, kicker: str = "") -> None:
    add_text(slide, f"{number:02d}", 0.45, 0.25, 0.48, 0.35, size=11, bold=True, color=BLUE)
    add_text(slide, title, 1.0, 0.2, 11.8, 0.58, size=25, bold=True, font=TITLE_FONT)
    if kicker:
        add_text(slide, kicker, 1.02, 0.76, 11.5, 0.34, size=11.5, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(1.14), Inches(12.4), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, index: int) -> None:
    add_text(slide, "자료: 「유전 알고리즘으로 알파 리서치하기」", 0.48, 7.13, 7.0, 0.2, size=8.5, color=MUTED)
    add_text(slide, str(index), 12.2, 7.1, 0.6, 0.22, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_note(slide, script: str) -> None:
    frame = slide.notes_slide.notes_text_frame
    frame.text = script


def add_arrow(slide, x1, y1, x2, y2, color=MUTED, width=2.0):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = PptPt(width)
    line.line.end_arrowhead = True
    return line


def add_bar(slide, label, value, max_value, x, y, w, *, color=BLUE, suffix="%"):
    add_text(slide, label, x, y - 0.02, 1.2, 0.33, size=13, bold=True)
    add_rect(slide, x + 1.15, y, w, 0.34, fill=GRAY_CARD, line=GRAY_CARD, radius=False)
    fill_w = max(0.02, w * value / max_value)
    add_rect(slide, x + 1.15, y, fill_w, 0.34, fill=color, line=color, radius=False)
    add_text(slide, f"{value:.2f}{suffix}", x + 1.15 + w + 0.12, y - 0.02, 1.1, 0.33, size=14, bold=True, color=color)


def add_picture_contain(slide, path: Path, x, y, w, h):
    with Image.open(path) as image:
        iw, ih = image.size
    ratio = min(w / iw, h / ih)
    pw, ph = iw * ratio, ih * ratio
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def make_slides(prs: Presentation, identity: str, assets: Path, scripts: list[dict]) -> None:
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, INK)
    add_rect(slide, 0.55, 0.55, 0.12, 5.9, fill=ORANGE, line=ORANGE, radius=False)
    add_text(slide, "유전 알고리즘으로\n알파 리서치하기", 1.0, 1.0, 9.8, 1.55, size=34, bold=True, color=WHITE, font=TITLE_FONT)
    add_text(slide, "1축 변이는 다축 변이보다 효과적인가", 1.02, 2.78, 9.8, 0.55, size=21, bold=True, color=RGBColor(214, 226, 235))
    add_text(slide, "외부 통제 실험 → GenomicWQB 퀀트 알파 실증", 1.02, 3.55, 10.3, 0.45, size=16, color=RGBColor(190, 198, 205))
    add_text(slide, identity, 1.02, 5.55, 10.8, 0.35, size=11.5, color=RGBColor(190, 198, 205))
    add_text(slide, "10분 발표", 10.8, 0.65, 1.7, 0.42, size=12, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)
    add_note(slide, scripts[0]["script"])

    # 2. Research question
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, 2, "한 번에 몇 축을 바꿀 것인가", "변이 폭은 탐색 예산을 어디에 쓸지 정한다")
    add_card(slide, 0.65, 1.55, 5.65, 3.25, "1축 변이", "가까이, 정확하게", "구조 보존 · 원인 추적 · 좋은 부모 주변", accent=BLUE, value_size=25)
    add_card(slide, 7.0, 1.55, 5.65, 3.25, "다축 변이", "멀리, 한 번에", "새 영역 탐색 · 정체 탈출 · 축 간 상호작용", accent=ORANGE, value_size=25)
    add_text(slide, "연구 질문", 0.8, 5.25, 1.3, 0.35, size=12, bold=True, color=MUTED)
    add_text(slide, "좋은 후보를 다듬을 때 1축 변이가 더 효과적인가?", 2.0, 5.12, 10.2, 0.55, size=24, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_footer(slide, 2)
    add_note(slide, scripts[1]["script"])

    # 3. Evidence order
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, 3, "검증 순서가 결론의 힘을 만든다", "이론을 통제 실험으로 확인하고 실제 알파에서 다시 검증")
    steps = [
        ("1", "이론", "국소 탐색은\n좁은 변이가 유리", BLUE_LIGHT, BLUE),
        ("2", "외부 실험", "PBO 3,000쌍\n변이 폭만 통제", BLUE_LIGHT, BLUE),
        ("3", "WQB 실증", "실평가 계보\n9,839쌍", ORANGE_LIGHT, ORANGE),
        ("4", "운용 원칙", "활용 1~2축\n정체 시 3축+", GRAY_CARD, INK),
    ]
    x_positions = [0.6, 3.8, 7.0, 10.2]
    for i, ((num, title, body, fill, accent), x) in enumerate(zip(steps, x_positions)):
        add_rect(slide, x, 1.8, 2.55, 3.35, fill=fill, line=accent)
        add_text(slide, num, x + 0.2, 2.02, 0.45, 0.42, size=16, bold=True, color=accent)
        add_text(slide, title, x + 0.2, 2.58, 2.15, 0.44, size=20, bold=True, font=TITLE_FONT)
        add_text(slide, body, x + 0.2, 3.35, 2.15, 0.95, size=15, color=INK, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_arrow(slide, x + 2.58, 3.45, x_positions[i + 1] - 0.08, 3.45, color=MUTED)
    add_text(slide, "서로 다른 자료에서 같은 방향이 나와야 강건한 결론이다", 1.3, 5.75, 10.7, 0.5, size=21, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_footer(slide, 3)
    add_note(slide, scripts[2]["script"])

    # 4. External design
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, 4, "외부 통제 실험", "IOHprofiler PBO에서 시작점·예산을 같게 두고 변이 폭만 바꿨다")
    add_card(slide, 0.6, 1.45, 2.8, 1.65, "공개 문제", "25종", "PBO 문제군", accent=BLUE)
    add_card(slide, 3.6, 1.45, 2.8, 1.65, "대응 비교", "3,000쌍", "같은 시작점", accent=BLUE)
    add_card(slide, 6.6, 1.45, 2.8, 1.65, "알고리즘 실행", "6,000회", "1축과 다축", accent=ORANGE)
    add_card(slide, 9.6, 1.45, 3.05, 1.65, "함수평가", "696만 회", "예산 20d", accent=ORANGE)
    add_rect(slide, 0.8, 3.6, 5.7, 1.8, fill=BLUE_LIGHT, line=BLUE)
    add_text(slide, "1축", 1.1, 3.9, 1.0, 0.5, size=26, bold=True, color=BLUE)
    add_text(slide, "매 평가마다 한 자리 변경", 2.25, 3.96, 3.8, 0.4, size=17, bold=True)
    add_rect(slide, 6.85, 3.6, 5.7, 1.8, fill=ORANGE_LIGHT, line=ORANGE)
    add_text(slide, "다축", 7.15, 3.9, 1.0, 0.5, size=26, bold=True, color=ORANGE)
    add_text(slide, "WQB의 2~14축 비율대로 변경", 8.3, 3.96, 3.85, 0.52, size=17, bold=True)
    add_text(slide, "핵심 지표: 최적점 도달률  |  보조 지표: 최적점까지 줄인 거리", 1.2, 5.85, 10.9, 0.45, size=18, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4)
    add_note(slide, scripts[3]["script"])

    # 5. External result
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, 5, "외부 결과에서 국소 수렴은 1축이 강했다", "최적점 도달률은 1축이 높지만 험준한 문제에서는 다축 점프가 필요")
    add_text(slide, "최적점 도달률", 0.8, 1.4, 4.8, 0.42, size=16, bold=True)
    add_bar(slide, "1축", 45.5, 50, 0.8, 2.0, 3.7, color=BLUE)
    add_bar(slide, "다축", 26.5, 50, 0.8, 2.75, 3.7, color=ORANGE)
    add_text(slide, "+19.0%p", 1.7, 3.55, 3.9, 0.65, size=34, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "95% CI [8.8, 29.1]", 1.7, 4.2, 3.9, 0.35, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 6.7, 1.55, 5.8, 1.55, fill=BLUE_LIGHT, line=BLUE)
    add_text(slide, "기본 문제", 7.0, 1.82, 1.6, 0.35, size=14, bold=True, color=BLUE)
    add_text(slide, "1축 도달률 +69.2%p", 8.55, 1.78, 3.45, 0.48, size=22, bold=True, color=BLUE)
    add_rect(slide, 6.7, 3.35, 5.8, 1.55, fill=ORANGE_LIGHT, line=ORANGE)
    add_text(slide, "험준·상호작용", 7.0, 3.62, 1.85, 0.35, size=14, bold=True, color=ORANGE)
    add_text(slide, "다축 이동 거리 +0.230", 8.75, 3.58, 3.3, 0.48, size=21, bold=True, color=ORANGE)
    add_text(slide, "이론적 토대", 0.9, 5.35, 1.4, 0.35, size=12, bold=True, color=MUTED)
    add_text(slide, "1축은 ‘좋은 후보를 다듬는 연산자’다. 장벽 탈출까지 맡기면 안 된다.", 2.2, 5.2, 10.0, 0.7, size=23, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_footer(slide, 5)
    add_note(slide, scripts[4]["script"])

    # 6. Internal data
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, 6, "퀀트 알파에서도 같은 현상이 나오는가", "GenomicWQB 운영 기록을 실제 평가가 일어난 부모·자식 쌍으로 정리")
    stages = [
        ("전체 알파", "39,506건", BLUE),
        ("Sharpe 연결", "12,168쌍", BLUE),
        ("실평가 계보", "9,839쌍", GREEN),
    ]
    xs = [0.75, 4.65, 8.55]
    for i, ((label, value, accent), x) in enumerate(zip(stages, xs)):
        add_rect(slide, x, 1.55, 3.15, 1.55, fill=WHITE, line=accent)
        add_text(slide, label, x + 0.2, 1.82, 2.75, 0.3, size=13, bold=True, color=MUTED)
        add_text(slide, value, x + 0.2, 2.18, 2.75, 0.5, size=27, bold=True, color=accent, align=PP_ALIGN.CENTER)
        if i < 2:
            add_arrow(slide, x + 3.2, 2.33, xs[i + 1] - 0.08, 2.33)
    add_card(slide, 2.2, 3.7, 3.7, 1.55, "1축", "3,402쌍", "genes_changed = 1", accent=BLUE)
    add_card(slide, 7.25, 3.7, 3.7, 1.55, "다축", "6,437쌍", "genes_changed ≥ 2", accent=ORANGE)
    add_text(slide, "검증 3종: 단순 비교 · 운영 조건 보정 · 같은 부모의 자식 비교", 1.2, 5.8, 10.9, 0.45, size=18, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 6)
    add_note(slide, scripts[5]["script"])

    # 7. WQB main result
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, 7, "WQB에서도 1축이 부모를 더 자주 개선했다", "부모 대비 Sharpe 개선률")
    add_bar(slide, "1축", 16.05, 20, 0.7, 1.7, 4.1, color=BLUE)
    add_bar(slide, "다축", 11.96, 20, 0.7, 2.5, 4.1, color=ORANGE)
    add_text(slide, "+4.09%p", 1.6, 3.35, 4.3, 0.7, size=38, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "95% CI [1.64, 6.49]", 1.6, 4.05, 4.3, 0.35, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    add_card(slide, 7.0, 1.55, 2.6, 1.65, "운영 조건 보정", "+4.77%p", "95% CI [2.53, 7.01]", accent=BLUE, value_size=24)
    add_card(slide, 9.9, 1.55, 2.6, 1.65, "같은 부모 비교", "+4.51%p", "95% CI [2.14, 6.89]", accent=BLUE, value_size=24)
    effect_img = assets / "report-img-004.jpg"
    add_picture_contain(slide, effect_img, 6.85, 3.5, 5.8, 2.1)
    add_text(slide, "세 계산이 모두 같은 방향", 7.2, 5.85, 5.1, 0.4, size=17, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_footer(slide, 7)
    add_note(slide, scripts[6]["script"])

    # 8. Dose response
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, 8, "정확한 경계는 1축 대 다축이 아니다", "1~2축의 좁은 변이와 3축 이상의 넓은 변이 사이에서 성과가 갈렸다")
    dose_img = assets / "report-img-002.jpg"
    add_picture_contain(slide, dose_img, 0.45, 1.35, 9.15, 4.55)
    add_card(slide, 9.75, 1.6, 2.8, 1.35, "1축", "16.05%", "부모 개선률", accent=BLUE, value_size=24)
    add_card(slide, 9.75, 3.15, 2.8, 1.35, "2축", "16.73%", "1축과 오차 범위 겹침", accent=BLUE, value_size=24)
    add_card(slide, 9.75, 4.7, 2.8, 1.35, "5축 이상", "약 6%", "개선률 급락", accent=ORANGE, value_size=24)
    add_text(slide, "주의: 모든 구간의 ΔSharpe 중앙값은 음수 → 절대 성공이 아니라 상대 우위", 0.8, 6.2, 11.7, 0.42, size=15.5, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_footer(slide, 8)
    add_note(slide, scripts[7]["script"])

    # 9. Robust conclusion and boundary
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, 9, "강건하지만 범위가 있는 결론", "자료와 계산을 바꿔도 국소 개선 우위는 유지됐지만 최종 기준 도달은 늘지 않았다")
    add_card(slide, 0.65, 1.5, 3.75, 1.6, "단순 비교", "+4.09%p", "부모 군집 부트스트랩", accent=BLUE)
    add_card(slide, 4.78, 1.5, 3.75, 1.6, "조건 보정", "+4.77%p", "시점·경로·시장 설정 통제", accent=BLUE)
    add_card(slide, 8.91, 1.5, 3.75, 1.6, "같은 부모", "+4.51%p", "부모 품질 차이 축소", accent=BLUE)
    add_rect(slide, 0.8, 3.7, 5.75, 1.7, fill=BLUE_LIGHT, line=BLUE)
    add_text(slide, "자료가 지지한 결론", 1.1, 3.98, 2.2, 0.35, size=14, bold=True, color=BLUE)
    add_text(slide, "국소 탐색에서 1~2축은 유효하다", 1.25, 4.45, 4.9, 0.44, size=21, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 6.8, 3.7, 5.75, 1.7, fill=ORANGE_LIGHT, line=ORANGE)
    add_text(slide, "자료가 지지하지 않은 결론", 7.1, 3.98, 2.8, 0.35, size=14, bold=True, color=ORANGE)
    add_text(slide, "1축이 항상 최고이거나 최종 성공을 보장한다", 7.2, 4.38, 5.0, 0.64, size=19, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "S≥1.58 신규 도달: 1축 1.56% · 다축 1.48% → 뚜렷한 차이 없음", 1.0, 5.9, 11.3, 0.4, size=16, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 9)
    add_note(slide, scripts[8]["script"])

    # 10. Operating rule
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, 10, "운용 규칙은 활용과 탈출을 나눈다", "1~2축은 기본 연산자, 3축 이상은 정체를 벗어나는 조건부 연산자")
    add_rect(slide, 0.7, 1.65, 2.7, 1.45, fill=BLUE_LIGHT, line=BLUE)
    add_text(slide, "유망한 부모", 1.0, 1.95, 2.1, 0.38, size=21, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, "최근 개선 지속", 1.0, 2.45, 2.1, 0.28, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_arrow(slide, 3.42, 2.38, 4.28, 2.38, color=BLUE)
    add_rect(slide, 4.3, 1.65, 3.1, 1.45, fill=BLUE, line=BLUE)
    add_text(slide, "1~2축 우선", 4.65, 1.98, 2.4, 0.45, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "국소 활용", 4.8, 2.48, 2.1, 0.28, size=12, color=RGBColor(220, 232, 240), align=PP_ALIGN.CENTER)
    add_arrow(slide, 7.42, 2.38, 8.28, 2.38, color=MUTED)
    add_rect(slide, 8.3, 1.65, 2.2, 1.45, fill=WHITE, line=MUTED)
    add_text(slide, "정체?", 8.7, 2.02, 1.4, 0.42, size=24, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "연속 실패", 8.72, 2.48, 1.35, 0.28, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_arrow(slide, 10.52, 2.38, 11.0, 2.38, color=ORANGE)
    add_rect(slide, 11.02, 1.65, 1.7, 1.45, fill=ORANGE, line=ORANGE)
    add_text(slide, "3축+", 11.25, 1.98, 1.25, 0.45, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "장벽 탈출", 11.2, 2.48, 1.35, 0.28, size=12, color=RGBColor(250, 233, 215), align=PP_ALIGN.CENTER)
    add_rect(slide, 1.0, 4.05, 11.3, 1.55, fill=WHITE, line=LINE)
    add_text(slide, "다음 검증", 1.3, 4.33, 1.5, 0.33, size=14, bold=True, color=MUTED)
    add_text(slide, "같은 부모를 k=1 · k=2 · k≥3으로 무작위 배정", 2.65, 4.28, 6.9, 0.44, size=21, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "방법마다 최소 약 1,176쌍", 9.55, 4.3, 2.35, 0.44, size=17, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "관측 결론 → 인과 결론", 3.8, 5.95, 5.8, 0.5, size=23, bold=True, font=TITLE_FONT, color=GREEN, align=PP_ALIGN.CENTER)
    add_footer(slide, 10)
    add_note(slide, scripts[9]["script"])

    # 11. Final
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, INK)
    add_text(slide, "결론", 0.8, 0.75, 2.0, 0.45, size=16, bold=True, color=ORANGE)
    add_text(slide, "1축 변이는 퀀트 알파의\n국소 탐색에서 유효하다", 0.8, 1.35, 11.1, 1.5, size=35, bold=True, color=WHITE, font=TITLE_FONT)
    takeaways = [
        "외부 통제 실험과 WQB 계보가 같은 방향을 보였다",
        "정확한 운용 경계는 1~2축 대 3축 이상이다",
        "다축은 버릴 대상이 아니라 정체 탈출용 연산자다",
    ]
    for i, text in enumerate(takeaways, 1):
        y = 3.35 + (i - 1) * 0.72
        add_text(slide, f"{i}", 1.0, y, 0.4, 0.36, size=15, bold=True, color=ORANGE)
        add_text(slide, text, 1.55, y - 0.02, 10.3, 0.44, size=19, color=RGBColor(225, 229, 232))
    add_text(slide, "Q&A", 10.4, 6.35, 1.8, 0.48, size=23, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)
    add_note(slide, scripts[10]["script"])


def set_cell_text(cell, text, *, bold=False, size=9, color=DocxRGBColor(0, 0, 0)):
    cell.text = ""
    paragraph = cell.paragraphs[0]
    paragraph.paragraph_format.line_spacing = 1.0
    paragraph.paragraph_format.space_before = Pt(0)
    paragraph.paragraph_format.space_after = Pt(0)
    run = paragraph.add_run(text)
    run.font.name = BODY_FONT
    run._element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), BODY_FONT)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def shade_docx_cell(cell, fill: str) -> None:
    properties = cell._tc.get_or_add_tcPr()
    shading = properties.find(qn("w:shd"))
    if shading is None:
        shading = OxmlElement("w:shd")
        properties.append(shading)
    shading.set(qn("w:fill"), fill)


def make_script_docx(path: Path, identity: str, scripts: list[dict], report_hash: str) -> None:
    doc = Document()
    section = doc.sections[0]
    section.page_width = Mm(210)
    section.page_height = Mm(297)
    section.top_margin = Mm(18)
    section.bottom_margin = Mm(17)
    section.left_margin = Mm(20)
    section.right_margin = Mm(20)

    normal = doc.styles["Normal"]
    normal.font.name = BODY_FONT
    normal._element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), BODY_FONT)
    normal.font.size = Pt(10.3)
    normal.paragraph_format.line_spacing = 1.45
    normal.paragraph_format.space_after = Pt(5)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.paragraph_format.space_before = Pt(36)
    run = title.add_run("유전 알고리즘으로 알파 리서치하기")
    run.font.name = TITLE_FONT
    run._element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), TITLE_FONT)
    run.font.size = Pt(22)
    run.font.bold = True
    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run("10분 발표 대본과 슬라이드 큐시트")
    run.font.name = TITLE_FONT
    run._element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), TITLE_FONT)
    run.font.size = Pt(15)
    run.font.bold = True
    author = doc.add_paragraph(identity)
    author.alignment = WD_ALIGN_PARAGRAPH.CENTER
    author.paragraph_format.space_before = Pt(24)
    meta = doc.add_paragraph(f"총 11장 · 10분 00초 · 원본 PDF SHA-256 {report_hash}")
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("한눈에 보는 큐시트", level=1)
    table = doc.add_table(rows=1, cols=4)
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    headers = ["슬라이드", "구간", "시간", "한 문장 메시지"]
    for i, header in enumerate(headers):
        shade_docx_cell(table.rows[0].cells[i], "D9E6EF")
        set_cell_text(table.rows[0].cells[i], header, bold=True, size=9)
    elapsed = 0
    for item in scripts:
        start = elapsed
        elapsed += item["seconds"]
        row = table.add_row()
        values = [
            str(item["slide"]),
            f"{start // 60}:{start % 60:02d}–{elapsed // 60}:{elapsed % 60:02d}",
            f"{item['seconds']}초",
            item["message"],
        ]
        for i, value in enumerate(values):
            row.cells[i].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            set_cell_text(row.cells[i], value, size=8.6)

    doc.add_page_break()
    doc.add_heading("슬라이드별 발표 대본", level=1)
    elapsed = 0
    for index, item in enumerate(scripts, 1):
        start = elapsed
        elapsed += item["seconds"]
        heading = doc.add_heading(
            f"슬라이드 {item['slide']}. {item['title']}  "
            f"({start // 60}:{start % 60:02d}–{elapsed // 60}:{elapsed % 60:02d})",
            level=2,
        )
        heading.paragraph_format.keep_with_next = True
        message = doc.add_paragraph()
        message.paragraph_format.left_indent = Mm(4)
        message.paragraph_format.right_indent = Mm(4)
        run = message.add_run("화면 메시지  " + item["message"])
        run.font.bold = True
        run.font.color.rgb = DocxRGBColor(42, 95, 137)
        script = doc.add_paragraph(item["script"])
        script.paragraph_format.first_line_indent = Mm(4)
        transition = doc.add_paragraph("전환  " + item["transition"])
        transition.paragraph_format.left_indent = Mm(4)
        transition.runs[0].font.italic = True
        transition.runs[0].font.color.rgb = DocxRGBColor(92, 101, 111)
        if index in {3, 6, 9}:
            doc.add_page_break()

    doc.add_heading("발표 직전 점검", level=1)
    checks = [
        "수치는 슬라이드에 적힌 값만 읽고 소수점 자릿수를 늘리지 않는다.",
        "‘1축이 항상 우월하다’고 말하지 않는다. 국소 탐색이라는 범위를 붙인다.",
        "8번 슬라이드에서는 모든 ΔSharpe 중앙값이 음수라는 점을 반드시 언급한다.",
        "10번 슬라이드에서 현재 결론이 관측 결론이며 후속 무작위 실험이 필요하다고 닫는다.",
        "질문이 나오면 외부 실험은 작동 원리, WQB는 퀀트 알파의 적용 가능성을 맡는다고 답한다.",
    ]
    for text in checks:
        paragraph = doc.add_paragraph(style="List Bullet")
        paragraph.add_run(text)

    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(path)


def slide_scripts() -> list[dict]:
    return [
        {
            "slide": 1, "title": "표지", "seconds": 30,
            "message": "외부 통제 실험과 실제 알파 기록으로 1축 변이를 검증했다.",
            "script": "안녕하세요. 유전 알고리즘으로 알파를 탐색할 때 한 번에 축 하나만 바꾸는 방식이 실제로 유리한지 발표하겠습니다. 공개 최적화 문제에서 작동 원리를 먼저 확인하고 GenomicWQB의 실제 알파 계보에서 같은 현상이 나타나는지 검정했습니다. 발표의 초점은 1축이 언제 유효하고 언제 한계가 생기는지입니다.",
            "transition": "먼저 왜 변이 폭이 중요한지 보겠습니다.",
        },
        {
            "slide": 2, "title": "연구 질문", "seconds": 50,
            "message": "변이 폭은 가까운 곳을 다듬을지, 멀리 탐색할지를 정한다.",
            "script": "알파 하나에는 데이터필드, 연산자, 룩백, 중립화, 감쇠처럼 여러 축이 들어갑니다. 1축 변이는 부모 구조를 거의 보존해서 성적 변화의 원인을 찾기 쉽고 이미 좋은 후보 주변을 세밀하게 볼 수 있습니다. 다축 변이는 한 번에 더 멀리 가기 때문에 정체된 영역을 벗어나는 데 유리합니다. 그래서 질문을 이렇게 좁혔습니다. 좋은 부모를 다듬는 단계라면 1축 변이가 다축보다 효과적인가입니다.",
            "transition": "이 질문을 한 자료만으로 판단하지 않았습니다.",
        },
        {
            "slide": 3, "title": "검증 순서", "seconds": 50,
            "message": "이론 → 외부 통제 실험 → WQB 실증의 순서로 증거를 쌓았다.",
            "script": "검증은 네 단계로 이어집니다. 한 자리씩 고쳐도 성적이 좋아지는 문제라면 좁은 변이가 유리하다는 예측에서 출발했습니다. 이 예측을 PBO 공개 문제에서 먼저 확인했습니다. 그다음 WQB 운영 기록을 살펴 실제 퀀트 알파에서도 같은 현상이 나오는지 검정했습니다. 두 결과를 합친 뒤에는 운용 규칙을 정리했습니다. 외부 실험은 원리를, WQB는 현실 적용을 맡습니다. 서로 다른 자료가 같은 방향을 보여야 결론도 단단해집니다.",
            "transition": "외부 실험의 조건부터 말씀드리겠습니다.",
        },
        {
            "slide": 4, "title": "외부 통제 실험", "seconds": 50,
            "message": "같은 시작점과 예산에서 바꾸는 축 수만 달리했다.",
            "script": "외부 실험에는 IOHprofiler의 PBO 25개 문제를 사용했습니다. 문제 크기 두 종류와 인스턴스 세 종류를 각각 20회 반복해 3,000개의 대응쌍을 만들었습니다. 두 알고리즘은 같은 시작점에서 출발하고 평가예산도 20d로 같습니다. 1축은 매번 한 자리만 바꿉니다. 다축은 WQB에서 실제로 관측된 2축에서 14축 사이의 비율을 따릅니다. 전체 실행은 6,000회, 함수평가는 696만 회입니다.",
            "transition": "변이 폭만 바꿨을 때 결과는 분명한 경계로 나타났습니다.",
        },
        {
            "slide": 5, "title": "외부 실험 결과", "seconds": 70,
            "message": "1축은 최적점에 잘 모였고, 다축은 험준한 장벽에서 더 멀리 갔다.",
            "script": "가장 중요한 결과입니다. 최적점 도달률은 1축 45.5%, 다축 26.5%로 19.0퍼센트포인트 벌어졌습니다. 신뢰구간도 0보다 큽니다. 기본 문제만 따로 보면 1축의 도달률 우위는 69.2퍼센트포인트였습니다. 한 자리씩 고쳐도 되는 문제에서는 1축이 좋은 답에 더 안정적으로 모였습니다. 다만 험준하거나 축끼리 강하게 얽힌 문제에서는 다축의 이동 거리가 0.230만큼 길었습니다. 외부 실험이 보여 준 역할 구분은 분명합니다. 1축은 좋은 후보를 다듬고 다축은 장벽을 넘습니다.",
            "transition": "이제 이 원리가 퀀트 알파에서도 나타나는지 보겠습니다.",
        },
        {
            "slide": 6, "title": "GenomicWQB 자료", "seconds": 50,
            "message": "실제로 새 평가를 돌린 부모·자식 9,839쌍을 비교했다.",
            "script": "WQB에는 알파 39,506건이 있습니다. 부모와 자식의 Sharpe를 연결한 뒤 무변경, 캐시 재사용, 완전히 같은 설정을 제외했습니다. 최종 표본은 실제 평가가 일어난 부모·자식 9,839쌍입니다. 1축은 3,402쌍, 다축은 6,437쌍입니다. 결과가 특정 계산에만 의존하지 않도록 단순 비교, 운영 조건을 맞춘 비교, 같은 부모에서 나온 자식끼리의 비교를 함께 사용했습니다.",
            "transition": "세 방법의 결과는 거의 같은 크기로 모였습니다.",
        },
        {
            "slide": 7, "title": "WQB 핵심 결과", "seconds": 75,
            "message": "1축의 부모 대비 개선률은 다축보다 4~5%p 높았다.",
            "script": "부모보다 Sharpe가 좋아진 비율은 1축 16.05%, 다축 11.96%입니다. 단순 차이는 4.09퍼센트포인트이고 95% 신뢰구간은 1.64에서 6.49입니다. 평가 시점과 생성 방식, 리전, 유니버스 같은 운영 조건을 맞추면 차이는 4.77퍼센트포인트였습니다. 같은 부모에서 나온 자식끼리만 비교해도 4.51퍼센트포인트입니다. 세 계산이 같은 방향과 비슷한 크기를 보였습니다. 외부 실험에서 확인한 국소 탐색의 1축 우위가 실제 퀀트 알파 계보에서도 다시 나타났습니다.",
            "transition": "다만 정확한 경계는 1축과 나머지 전체를 가르는 선이 아니었습니다.",
        },
        {
            "slide": 8, "title": "변이 폭별 결과", "seconds": 65,
            "message": "실제 경계는 1~2축의 좁은 변이와 3축 이상의 넓은 변이다.",
            "script": "축 수를 나눠 보면 1축 개선률은 16.05%, 2축은 16.73%로 거의 같습니다. 3축부터 13.13%로 내려가고 5축 이상은 약 6%까지 떨어집니다. 1축만 혼자 좋다고 해석해서는 안 됩니다. 자료가 보여 주는 경계는 1축과 2축의 좁은 변이, 그리고 3축 이상의 넓은 변이 사이입니다. 주의할 점도 있습니다. 모든 구간에서 Sharpe 변화량의 중앙값은 음수입니다. 1축도 평균적으로는 성공하지 못했습니다. 다축보다 실패가 적고 하락폭이 작았다는 상대 우위입니다.",
            "transition": "이 구분을 지키면 결론의 범위도 명확해집니다.",
        },
        {
            "slide": 9, "title": "강건성과 한계", "seconds": 55,
            "message": "국소 개선 우위는 강건하지만 최종 성공과 보편적 우위는 확인되지 않았다.",
            "script": "단순 비교는 4.09, 조건을 맞춘 비교는 4.77, 같은 부모 비교는 4.51퍼센트포인트입니다. 계산법을 바꿔도 결과는 흔들리지 않았습니다. 하지만 Sharpe 1.58을 새로 넘은 비율은 1축 1.56%, 다축 1.48%로 뚜렷한 차이가 없습니다. 자료가 지지한 결론은 국소 탐색에서 1축과 2축이 유효하다는 데까지입니다. 1축이 언제나 최고이거나 최종 제출 성공을 보장한다는 결론은 자료가 받쳐 주지 않습니다.",
            "transition": "이 결론을 실제 탐색 정책으로 옮기면 다음과 같습니다.",
        },
        {
            "slide": 10, "title": "운용 규칙과 후속 검증", "seconds": 70,
            "message": "유망한 부모에는 1~2축, 정체가 확인되면 3축 이상을 투입한다.",
            "script": "운용 규칙은 상태에 따라 나눕니다. 유망한 부모가 최근에도 개선되고 있다면 1축과 2축을 먼저 사용합니다. 좁은 변이가 연속으로 실패하거나 여러 축을 함께 바꿔야 풀리는 구조로 보이면 3축 이상을 조건부로 투입합니다. 다축은 없애지 않고 장벽 탈출 역할로 옮깁니다. 다음 검증에서는 같은 부모를 1축, 2축, 3축 이상으로 무작위 배정해야 합니다. 4퍼센트포인트 차이를 찾으려면 방법마다 최소 약 1,176쌍이 필요합니다. 이 실험이 관측 결론을 인과 결론으로 바꾸는 단계입니다.",
            "transition": "마지막으로 세 문장만 남기겠습니다.",
        },
        {
            "slide": 11, "title": "결론", "seconds": 35,
            "message": "1축은 국소 활용, 다축은 조건부 탈출에 배치한다.",
            "script": "결론은 세 가지입니다. 외부 통제 실험과 WQB 계보가 같은 방향을 보였습니다. 퀀트 알파의 국소 탐색에서는 1축과 2축의 좁은 변이가 유효했습니다. 다축도 버릴 필요는 없습니다. 정체를 벗어날 때 쓰는 조건부 연산자로 두면 됩니다. 활용과 탈출의 역할을 나누는 것이 이번 연구가 제안하는 운용 원칙입니다. 감사합니다. 질문 받겠습니다.",
            "transition": "질의응답으로 넘어간다.",
        },
    ]


def write_audit(path: Path, report_hash: str, scripts: list[dict]) -> None:
    total_seconds = sum(item["seconds"] for item in scripts)
    total_chars = sum(len(item["script"]) for item in scripts)
    lines = [
        f"source_pdf_sha256={report_hash}",
        f"slides={len(scripts)}",
        f"total_seconds={total_seconds}",
        f"script_characters={total_chars}",
    ]
    elapsed = 0
    for item in scripts:
        start = elapsed
        elapsed += item["seconds"]
        lines.append(
            f"slide_{item['slide']:02d}={start // 60}:{start % 60:02d}-"
            f"{elapsed // 60}:{elapsed % 60:02d}|{item['title']}|{len(item['script'])}chars"
        )
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> None:
    args = parse_args()
    identity, _, report_hash = load_report(args.source_pdf)
    scripts = slide_scripts()
    if sum(item["seconds"] for item in scripts) != 600:
        raise RuntimeError("presentation timing must total exactly 600 seconds")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    make_slides(prs, identity, args.asset_dir, scripts)
    args.pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.pptx)

    make_script_docx(args.script_docx, identity, scripts, report_hash)
    write_audit(args.audit_text, report_hash, scripts)
    print(args.pptx)
    print(args.script_docx)


if __name__ == "__main__":
    main()
